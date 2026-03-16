#!/usr/bin/env python3
"""Generate Clone-Xs Overview PowerPoint from HTML slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colours ──
CRIMSON = RGBColor(0xDC, 0x14, 0x3C)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_TEXT = RGBColor(0x44, 0x44, 0x44)
SUBTITLE_CLR = RGBColor(0x55, 0x55, 0x55)
TABLE_TEXT = RGBColor(0x11, 0x11, 0x11)
CARD_BG = RGBColor(0xF5, 0xF5, 0xF5)
CARD_BORDER = RGBColor(0xE0, 0xE0, 0xE0)
ALT_ROW = RGBColor(0xF9, 0xF9, 0xF9)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
DARK_CARD = RGBColor(0x16, 0x21, 0x3E)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
INFO = RGBColor(0x29, 0x80, 0xB9)
WARN = RGBColor(0xF3, 0x9C, 0x12)

FONT = "Roboto Light"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ──
def _set_font(run, size=14, bold=False, color=BLACK, font_name=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    ea = rPr.find("a:ea", nsmap)
    if ea is None:
        ea = etree.SubElement(rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface", font_name)


def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text="", size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size, bold, color)
    return txBox


def add_slide_title(slide, title, left=Inches(0.7), top=Inches(0.5)):
    add_textbox(slide, left, top, Inches(10), Inches(0.5), title, size=24, bold=True, color=BLACK)


def add_separator(slide, y=Inches(1.05)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(11.9), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CRIMSON
    shape.line.fill.background()


def add_subtitle(slide, text, y=Inches(1.2)):
    add_textbox(slide, Inches(0.7), y, Inches(11), Inches(0.4), text, size=14, color=SUBTITLE_CLR)


def add_slide_number(slide, text):
    add_textbox(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.3), text, size=9, color=SUBTITLE_CLR, align=PP_ALIGN.RIGHT)


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_accent_bar(slide, left, top, height, color=CRIMSON):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_card_with_text(slide, left, top, width, height, title, body, accent=True, badge_text=None, badge_color=CRIMSON):
    add_rounded_rect(slide, left, top, width, height, CARD_BG, CARD_BORDER)
    if accent:
        add_accent_bar(slide, left, top + Inches(0.05), height - Inches(0.1))

    y_offset = top + Inches(0.12)
    if badge_text:
        bg = add_rounded_rect(slide, left + Inches(0.25), y_offset, Inches(0.8), Inches(0.22), badge_color)
        add_textbox(slide, left + Inches(0.25), y_offset - Pt(1), Inches(0.8), Inches(0.22),
                    badge_text, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y_offset += Inches(0.28)

    add_textbox(slide, left + Inches(0.25), y_offset, width - Inches(0.4), Inches(0.25),
                title, size=12, bold=True, color=BLACK)
    add_textbox(slide, left + Inches(0.25), y_offset + Inches(0.25), width - Inches(0.4), height - Inches(0.6),
                body, size=10, color=BODY_TEXT)


def add_kpi_box(slide, left, top, width, height, value, label, bg_color=CRIMSON):
    add_rounded_rect(slide, left, top, width, height, bg_color)
    add_textbox(slide, left, top + Inches(0.15), width, Inches(0.45), value,
                size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + height - Inches(0.35), width, Inches(0.25), label,
                size=10, color=WHITE, align=PP_ALIGN.CENTER)


def add_phase_card(slide, left, top, width, height, header_text, body_lines):
    # Header bar
    hdr_h = Inches(0.35)
    add_rounded_rect(slide, left, top, width, hdr_h + Inches(0.02), CRIMSON)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.03), width - Inches(0.3), hdr_h,
                header_text, size=10, bold=True, color=WHITE)
    # Body
    add_rounded_rect(slide, left, top + hdr_h - Inches(0.02), width, height - hdr_h + Inches(0.02), CARD_BG, CARD_BORDER)
    body = "\n".join(f"• {line}" for line in body_lines)
    add_textbox(slide, left + Inches(0.15), top + hdr_h + Inches(0.08), width - Inches(0.3), height - hdr_h - Inches(0.1),
                body, size=10, color=BODY_TEXT)


def add_highlight_box(slide, left, top, width, height, title, body):
    add_rounded_rect(slide, left, top, width, height, CRIMSON)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.25),
                title, size=12, bold=True, color=WHITE)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.35), width - Inches(0.4), height - Inches(0.45),
                body, size=10, color=WHITE)


def add_table(slide, left, top, width, rows_data, col_widths=None):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(0.35 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)

            if r_idx == 0:  # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLACK
                _set_font(run, size=11, bold=True, color=WHITE)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else ALT_ROW
                _set_font(run, size=10, color=TABLE_TEXT)
                if c_idx == 0:
                    run.font.bold = True

    return table_shape


# ════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, DARK_BG)
add_textbox(s, Inches(1), Inches(2.0), Inches(11), Inches(1), "Clone-Xs",
            size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(3.0), Inches(11), Inches(0.5),
            "Enterprise Unity Catalog Cloning for Databricks",
            size=20, color=CRIMSON, align=PP_ALIGN.CENTER)
# Badge
add_rounded_rect(s, Inches(5.4), Inches(3.8), Inches(2.5), Inches(0.35), CRIMSON)
add_textbox(s, Inches(5.4), Inches(3.8), Inches(2.5), Inches(0.35),
            "Open Source — MIT License", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
            "CLI  •  Web UI  •  REST API  •  Serverless  •  Notebooks",
            size=13, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)
add_slide_number(s, "1 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "The Problem")
add_separator(s)
add_subtitle(s, "Cloning Unity Catalog catalogs today is manual, fragile, and incomplete")

cards = [
    ("Manual SQL Scripts", "Teams write custom SQL to clone tables one-by-one. No automation, no error handling, no rollback."),
    ("Missing Metadata", "Permissions, tags, properties, ownership, constraints, and comments are silently dropped."),
    ("No Validation", "No way to verify clone accuracy — row counts, checksums, and schema parity are unchecked."),
    ("No Audit Trail", "Nobody knows who cloned what, when, or whether it succeeded. Compliance gaps everywhere."),
]
col_w = Inches(5.7)
gap = Inches(0.3)
for i, (title, body) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + col * (col_w + gap)
    y = Inches(1.7) + row * Inches(1.3)
    add_card_with_text(s, x, y, col_w, Inches(1.1), title, body)

add_highlight_box(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.7),
                  "Result",
                  "Cloning a single catalog takes hours, fails silently, and leaves teams without confidence in their data.")
add_slide_number(s, "2 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 3: THE SOLUTION
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "The Solution: Clone-Xs")
add_separator(s)
add_subtitle(s, "Enterprise-grade cloning toolkit — clone, compare, sync, and manage Unity Catalog catalogs")

kpis = [("56", "CLI Commands", CRIMSON), ("88", "Python Modules", BLACK),
        ("31", "Web UI Pages", INFO), ("61+", "API Endpoints", SUCCESS)]
kpi_w = Inches(2.75)
for i, (val, lbl, clr) in enumerate(kpis):
    add_kpi_box(s, Inches(0.7) + i * (kpi_w + Inches(0.2)), Inches(1.7), kpi_w, Inches(1.0), val, lbl, clr)

sol_cards = [
    ("CLI", "Command Line", "56 commands for cloning, diffing, syncing, rollback, validation, profiling, and more.", CRIMSON),
    ("WEB", "Web Dashboard", "31-page React UI with dark mode, dynamic dropdowns, real-time logs, and guided wizards.", INFO),
    ("API", "REST API", "61+ FastAPI endpoints for automation, CI/CD integration, and programmatic control.", SUCCESS),
]
card_w = Inches(3.75)
for i, (badge, title, body, clr) in enumerate(sol_cards):
    add_card_with_text(s, Inches(0.7) + i * (card_w + Inches(0.2)), Inches(3.0), card_w, Inches(1.3),
                       title, body, badge_text=badge, badge_color=clr)
add_slide_number(s, "3 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 4: KEY FEATURES
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Key Features")
add_separator(s)
add_subtitle(s, "Everything you need to clone, manage, and govern Unity Catalog at scale")

features = [
    ("CLONE", "Deep & Shallow Clone", "Full data copy or metadata-only. Incremental sync and time-travel support.", CRIMSON),
    ("METADATA", "Full Metadata Copy", "Permissions, ownership, tags, properties, security policies, constraints, comments.", CRIMSON),
    ("VALIDATE", "Post-Clone Validation", "Row count + checksum verification with automatic rollback on mismatch.", SUCCESS),
    ("SERVERLESS", "Serverless Compute", "Run clones without a SQL warehouse — uploads wheel, submits notebook job automatically.", INFO),
    ("AUDIT", "Run Logs to Delta", "Every operation persists logs, metrics, and audit trail to Unity Catalog Delta tables.", BLACK),
    ("SYNC", "Incremental Sync", "Sync only changed tables using Delta version history — no full re-clone needed.", WARN),
]
col_w = Inches(5.7)
for i, (badge, title, body, clr) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + col * (col_w + Inches(0.3))
    y = Inches(1.7) + row * Inches(1.2)
    add_card_with_text(s, x, y, col_w, Inches(1.0), title, body, badge_text=badge, badge_color=clr)
add_slide_number(s, "4 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 5: ARCHITECTURE
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Architecture")
add_separator(s)
add_subtitle(s, "Three-tier design with shared Python core")

arch_boxes = [
    ("React UI", "31 pages\nVite + TanStack Query\nTailwind CSS 4\nshadcn/ui", DARK_BG),
    ("FastAPI Backend", "61+ endpoints\nJob queue + WebSockets\nReal-time log streaming\nPydantic v2", CRIMSON),
    ("Python Core", "88 modules\nDatabricks SDK\nUnity Catalog APIs\nDelta Lake ops", DARK_BG),
    ("Databricks", "Unity Catalog\nSQL Warehouses\nServerless Compute\nDelta Tables", CARD_BG),
]
box_w = Inches(2.6)
arrow_w = Inches(0.4)
total_w = 4 * box_w + 3 * arrow_w
start_x = Inches(0.7)
for i, (title, desc, bg) in enumerate(arch_boxes):
    x = start_x + i * (box_w + arrow_w)
    add_rounded_rect(s, x, Inches(1.7), box_w, Inches(1.6), bg, CARD_BORDER if bg == CARD_BG else None)
    txt_color = WHITE if bg != CARD_BG else BLACK
    add_textbox(s, x + Inches(0.1), Inches(1.8), box_w - Inches(0.2), Inches(0.3),
                title, size=12, bold=True, color=txt_color, align=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.1), Inches(2.15), box_w - Inches(0.2), Inches(1.0),
                desc, size=9, color=txt_color if bg != CARD_BG else BODY_TEXT, align=PP_ALIGN.CENTER)
    if i < 3:
        add_textbox(s, x + box_w, Inches(2.2), arrow_w, Inches(0.4), "→",
                    size=24, bold=True, color=CRIMSON, align=PP_ALIGN.CENTER)

# Bottom cards
bottom_cards = [
    ("CLI Layer", "• Python Click framework\n• 56 commands, YAML config\n• Same core as API"),
    ("Job Queue", "• Async job execution\n• WebSocket progress\n• Concurrent schema processing"),
    ("Audit Pipeline", "• Auto-create Delta tables\n• Every operation logged\n• Queryable via SQL"),
]
for i, (title, body) in enumerate(bottom_cards):
    x = Inches(0.7) + i * Inches(4.1)
    add_rounded_rect(s, x, Inches(3.6), Inches(3.8), Inches(1.3), CARD_BG, CARD_BORDER)
    add_textbox(s, x + Inches(0.15), Inches(3.7), Inches(3.5), Inches(0.3),
                title, size=11, bold=True, color=BLACK)
    add_textbox(s, x + Inches(0.15), Inches(4.0), Inches(3.5), Inches(0.8),
                body, size=9, color=BODY_TEXT)
add_slide_number(s, "5 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 6: WEB UI PAGES
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Web UI — 31 Pages")
add_separator(s)
add_subtitle(s, "Modern React dashboard with dark mode, real-time logs, and guided wizards")

table_data = [
    ["Category", "Pages", "Count"],
    ["Overview", "Dashboard, Audit Trail, Metrics", "3"],
    ["Operations", "Clone, Sync, Incremental Sync, Generate, Rollback, Templates, Schedule, Multi-Clone", "8"],
    ["Discovery", "Explorer, Diff & Compare, Config Diff, Lineage, Dependencies, Impact Analysis, Data Preview", "7"],
    ["Analysis", "Reports, PII Scanner, Schema Drift, Profiling, Cost Estimator, Compliance", "6"],
    ["Management", "Monitor, Preflight, Config, Settings, Warehouse, RBAC, Plugins", "7"],
]
add_table(s, Inches(0.7), Inches(1.7), Inches(11.9), table_data,
          col_widths=[Inches(2), Inches(8.5), Inches(1.4)])

kpis_ui = [("17", "Dynamic Dropdown Pages", CRIMSON), ("12", "Clone Templates", BLACK), ("4", "Clone Wizard Steps", SUCCESS)]
for i, (val, lbl, clr) in enumerate(kpis_ui):
    add_kpi_box(s, Inches(0.7) + i * Inches(4.1), Inches(5.2), Inches(3.8), Inches(0.9), val, lbl, clr)
add_slide_number(s, "6 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 7: CLONE WIZARD
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Clone Wizard — 4-Step Guided Flow")
add_separator(s)
add_subtitle(s, "Source → Options → Preview → Execute — with real-time progress and logs")

phases = [
    ("1. SOURCE", ["Select source catalog", "Select destination catalog", "Filter by schema/table", "Include/exclude patterns"]),
    ("2. OPTIONS", ["Deep vs Shallow clone", "Copy permissions, tags", "Enable validation", "Serverless toggle"]),
    ("3. PREVIEW", ["Review all settings", "See object counts", "Dry-run mode", "Preflight checks"]),
    ("4. EXECUTE", ["Real-time log stream", "Progress bar per schema", "Live status badges", "Databricks job link"]),
]
phase_w = Inches(2.8)
for i, (hdr, items) in enumerate(phases):
    add_phase_card(s, Inches(0.7) + i * (phase_w + Inches(0.2)), Inches(1.7), phase_w, Inches(1.8), hdr, items)

# Bottom cards
add_card_with_text(s, Inches(0.7), Inches(3.8), Inches(5.7), Inches(1.2),
                   "15 Clone Options",
                   "Permissions, ownership, tags, properties, security, constraints, comments, validation, checksum, rollback, report, checkpoint, force re-clone, verbose, serverless.")
add_card_with_text(s, Inches(6.8), Inches(3.8), Inches(5.8), Inches(1.2),
                   "12 Built-in Templates",
                   "Production Mirror, Dev Sandbox, DR Copy, Staging Refresh, Compliance Clone, Data Lake Copy, Schema-Only, Table Subset, Cross-Region, Audit Clone, Performance Test, Migration.")
add_slide_number(s, "7 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 8: SERVERLESS
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Serverless Compute")
add_separator(s)
add_subtitle(s, "Run clones on Databricks serverless — no SQL warehouse needed")

# Left column: How It Works
add_textbox(s, Inches(0.7), Inches(1.65), Inches(3), Inches(0.3), "HOW IT WORKS",
            size=11, bold=True, color=BLACK)
steps = [
    ("1. Package & Upload", "Builds Python wheel, uploads to UC Volume."),
    ("2. Generate Notebook", "Creates runner notebook with all clone options."),
    ("3. Submit Job", "Submits one-time serverless job. No cluster provisioning."),
    ("4. Stream Results", "Polls job status, streams output to UI in real-time."),
]
for i, (title, body) in enumerate(steps):
    add_card_with_text(s, Inches(0.7), Inches(2.0) + i * Inches(1.1), Inches(5.7), Inches(0.9), title, body)

# Right column: Benefits
add_textbox(s, Inches(6.8), Inches(1.65), Inches(3), Inches(0.3), "BENEFITS",
            size=11, bold=True, color=BLACK)
add_kpi_box(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(0.8), "$0", "No Warehouse Costs", CRIMSON)
add_kpi_box(s, Inches(6.8), Inches(2.95), Inches(5.8), Inches(0.8), "Auto", "Scales Automatically", BLACK)
add_kpi_box(s, Inches(6.8), Inches(3.9), Inches(5.8), Inches(0.8), "Fast", "No Cluster Start Wait", SUCCESS)

# CLI box
add_rounded_rect(s, Inches(6.8), Inches(4.9), Inches(5.8), Inches(1.3), CARD_BG, CARD_BORDER)
add_textbox(s, Inches(7.0), Inches(5.0), Inches(5.4), Inches(0.25), "CLI Usage", size=11, bold=True, color=BLACK)
add_rounded_rect(s, Inches(7.0), Inches(5.3), Inches(5.4), Inches(0.75), DARK_BG)
add_textbox(s, Inches(7.1), Inches(5.35), Inches(5.2), Inches(0.65),
            "clxs clone\n  --source edp_dev\n  --dest edp_dev_00\n  --serverless\n  --volume /Volumes/cat/schema/vol",
            size=9, color=SUCCESS)
add_slide_number(s, "8 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 9: AUDIT & LOGGING
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Audit Trail & Run Logs")
add_separator(s)
add_subtitle(s, "Every operation auto-persists to Unity Catalog Delta tables — queryable via SQL")

audit_cards = [
    ("logs.run_logs", "Full execution trace", ["Job ID, type, status", "Source → destination", "All log lines as ARRAY", "Result JSON", "Config (sanitized)", "Duration, timestamps"]),
    ("logs.clone_operations", "Compliance audit trail", ["Operation ID & type", "Who cloned what, when", "Tables/views/functions/volumes", "Success vs failure counts", "Summary JSON", "Error messages"]),
    ("metrics.clone_metrics", "Performance metrics", ["Throughput rates", "Success/failure rates", "Per-table timing", "Schema-level aggregates", "Trend analysis"]),
]
card_w = Inches(3.75)
for i, (hdr, sub, items) in enumerate(audit_cards):
    x = Inches(0.7) + i * (card_w + Inches(0.2))
    add_phase_card(s, x, Inches(1.7), card_w, Inches(2.8), hdr, [sub] + items)

add_highlight_box(s, Inches(0.7), Inches(4.8), Inches(11.9), Inches(0.7),
                  "Operations Logged",
                  "Clone, Sync, Incremental Sync, Validate, PII Scan, Preflight, Schema Evolve, Diff, Rollback, Terraform Generate — all write to Delta automatically.")
add_slide_number(s, "9 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 10: INCREMENTAL SYNC
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Incremental Sync")
add_separator(s)
add_subtitle(s, "Sync only changed tables using Delta version history — no full re-clone")

# Left
add_textbox(s, Inches(0.7), Inches(1.65), Inches(3), Inches(0.3), "HOW IT WORKS",
            size=11, bold=True, color=BLACK)
how_items = [
    "Compares Delta table versions between source and dest",
    "Identifies tables with new commits since last sync",
    "Syncs only changed tables — skips up-to-date ones",
    "Supports schema-level and table-level selection",
    "Parallel processing across schemas",
    "Runs on serverless or SQL warehouse",
]
add_rounded_rect(s, Inches(0.7), Inches(2.0), Inches(5.7), Inches(2.2), CARD_BG, CARD_BORDER)
body = "\n".join(f"→ {item}" for item in how_items)
add_textbox(s, Inches(0.9), Inches(2.1), Inches(5.3), Inches(2.0), body, size=10, color=BODY_TEXT)

add_rounded_rect(s, Inches(0.7), Inches(4.4), Inches(5.7), Inches(1.6), CARD_BG, CARD_BORDER)
add_textbox(s, Inches(0.9), Inches(4.5), Inches(5.3), Inches(0.25), "Selection Modes", size=11, bold=True, color=BLACK)
modes = "→ All schemas — scan entire catalog\n→ Selected schemas — pick specific schemas\n→ Selected tables — per-table checkboxes\n→ Select all / Clear — bulk operations"
add_textbox(s, Inches(0.9), Inches(4.8), Inches(5.3), Inches(1.1), modes, size=10, color=BODY_TEXT)

# Right
add_textbox(s, Inches(6.8), Inches(1.65), Inches(3), Inches(0.3), "UI FEATURES",
            size=11, bold=True, color=BLACK)
add_kpi_box(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(0.7), "Live", "Schema Results Stream In", CRIMSON)

add_rounded_rect(s, Inches(6.8), Inches(2.9), Inches(5.8), Inches(1.4), CARD_BG, CARD_BORDER)
add_textbox(s, Inches(7.0), Inches(3.0), Inches(5.4), Inches(0.25), "Per-Schema Status", size=11, bold=True, color=BLACK)
add_textbox(s, Inches(7.0), Inches(3.3), Inches(5.4), Inches(0.9),
            "→ Up to date — green badge\n→ N changed — amber badge with count\n→ Expandable table list with checkboxes\n→ Version history per table",
            size=10, color=BODY_TEXT)

add_rounded_rect(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(1.5), CARD_BG, CARD_BORDER)
add_textbox(s, Inches(7.0), Inches(4.6), Inches(5.4), Inches(0.25), "Sync Execution", size=11, bold=True, color=BLACK)
add_textbox(s, Inches(7.0), Inches(4.9), Inches(5.4), Inches(0.9),
            "→ Jobs submitted per schema\n→ Real-time progress tracking\n→ Live logs streaming\n→ Databricks job run links",
            size=10, color=BODY_TEXT)
add_slide_number(s, "10 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 11: SECURITY
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Security & Governance")
add_separator(s)
add_subtitle(s, "Enterprise security built in — credentials never stored server-side")

sec_cards = [
    ("CREDENTIALS", "Browser-Only Storage", "Databricks tokens stored in browser session only. Passed via headers — never persisted on server.", CRIMSON),
    ("AUDIT", "Sanitized Logging", "All tokens and secrets stripped from config before writing to audit Delta tables.", BLACK),
    ("RBAC", "Role-Based Access Control", "Define policies for who can clone which catalogs. Enforceable at API and UI level.", INFO),
    ("PII", "PII Scanner", "Scan catalogs for PII patterns (email, phone, SSN) before cloning to prevent data exposure.", SUCCESS),
]
col_w = Inches(5.7)
for i, (badge, title, body, clr) in enumerate(sec_cards):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + col * (col_w + Inches(0.3))
    y = Inches(1.7) + row * Inches(1.2)
    add_card_with_text(s, x, y, col_w, Inches(1.0), title, body, badge_text=badge, badge_color=clr)

# Bottom row
add_rounded_rect(s, Inches(0.7), Inches(4.3), Inches(5.7), Inches(1.4), CARD_BG, CARD_BORDER)
add_textbox(s, Inches(0.9), Inches(4.4), Inches(5.3), Inches(0.25), "Compliance Features", size=11, bold=True, color=BLACK)
add_textbox(s, Inches(0.9), Inches(4.7), Inches(5.3), Inches(0.9),
            "→ Full audit trail with Delta time travel\n→ Who cloned what, when, with what config\n→ Immutable operation logs\n→ Compliance dashboard in UI",
            size=10, color=BODY_TEXT)

add_rounded_rect(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(1.4), CARD_BG, CARD_BORDER)
add_textbox(s, Inches(7.0), Inches(4.4), Inches(5.4), Inches(0.25), "Safety Features", size=11, bold=True, color=BLACK)
add_textbox(s, Inches(7.0), Inches(4.7), Inches(5.4), Inches(0.9),
            "→ Preflight checks before clone\n→ Dry-run mode for all operations\n→ Auto-rollback on validation failure\n→ Checkpoint & restore support",
            size=10, color=BODY_TEXT)
add_slide_number(s, "11 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 12: TECH STACK
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Tech Stack")
add_separator(s)
add_subtitle(s, "Modern, production-ready stack with first-class Databricks integration")

tech_data = [
    ["Layer", "Technology", "Details"],
    ["Frontend", "React 19 + TypeScript", "Vite build, TanStack Query, Tailwind CSS 4, shadcn/ui, Recharts"],
    ["Backend", "Python + FastAPI", "Uvicorn, Pydantic v2, WebSockets, async job queue"],
    ["Core Engine", "Databricks SDK", "Unity Catalog REST APIs, SQL Statement Execution API"],
    ["CLI", "Python Click", "56 commands, YAML config profiles, rich output formatting"],
    ["Storage", "Delta Lake", "Audit logs, run history, metrics — all in Unity Catalog tables"],
    ["Docs", "Docusaurus", "Full documentation site with search, versioning"],
    ["Deployment", "Docker + Wheel", "Docker Compose for full stack, wheel for Databricks notebook/serverless"],
]
add_table(s, Inches(0.7), Inches(1.7), Inches(11.9), tech_data,
          col_widths=[Inches(2), Inches(3), Inches(6.9)])

add_rounded_rect(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.7), CARD_BG, CARD_BORDER)
add_textbox(s, Inches(0.9), Inches(5.65), Inches(11.5), Inches(0.25), "Project Stats", size=12, bold=True, color=BLACK)
add_textbox(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.25),
            "88 Python modules  •  31 React pages  •  61+ API endpoints  •  56 CLI commands  •  12 templates  •  Full test suite",
            size=11, color=BODY_TEXT)
add_slide_number(s, "12 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 13: GETTING STARTED
# ════════════════════════════════════════════════════════════
s = add_slide()
add_slide_title(s, "Getting Started")
add_separator(s)
add_subtitle(s, "Up and running in 3 commands")

install_options = [
    ("OPTION 1: FROM SOURCE", "git clone github.com/viral0216/clone-xs\ncd clone-xs\npip install -e \".[dev]\"\nmake web-start", "Frontend on :3000 • API on :8000"),
    ("OPTION 2: DOCKER", "docker-compose up --build\n\n# Open http://localhost:8000", "Full stack in one container"),
    ("OPTION 3: CLI ONLY", "pip install -e .\nclxs clone \\\n  --source my_catalog \\\n  --dest my_catalog_clone", "No server needed for CLI usage"),
]
card_w = Inches(3.75)
for i, (hdr, code, outcome) in enumerate(install_options):
    x = Inches(0.7) + i * (card_w + Inches(0.2))
    # Header
    hdr_h = Inches(0.35)
    add_rounded_rect(s, x, Inches(1.7), card_w, hdr_h, CRIMSON)
    add_textbox(s, x + Inches(0.15), Inches(1.73), card_w - Inches(0.3), hdr_h,
                hdr, size=10, bold=True, color=WHITE)
    # Code block
    add_rounded_rect(s, x, Inches(2.05), card_w, Inches(1.5), DARK_BG)
    add_textbox(s, x + Inches(0.15), Inches(2.1), card_w - Inches(0.3), Inches(1.4),
                code, size=9, color=SUCCESS)
    # Outcome
    add_rounded_rect(s, x, Inches(3.55), card_w, Inches(0.35), RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(s, x + Inches(0.1), Inches(3.55), card_w - Inches(0.2), Inches(0.35),
                outcome, size=9, color=RGBColor(0x2E, 0x7D, 0x32), align=PP_ALIGN.CENTER)

# Bottom: Required setup
add_rounded_rect(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(1.3), CARD_BG, CARD_BORDER)
add_textbox(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.25),
            "Required Setup After Install", size=12, bold=True, color=BLACK)
add_textbox(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.8),
            "→ Connect to Databricks — Enter workspace URL + PAT in Settings\n→ Select SQL Warehouse — Choose from dropdown (or use serverless)\n→ Initialize Audit Tables — Optional: set audit catalog for Delta logging",
            size=10, color=BODY_TEXT)
add_slide_number(s, "13 / 14")

# ════════════════════════════════════════════════════════════
# SLIDE 14: CLOSING
# ════════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, DARK_BG)
add_textbox(s, Inches(1), Inches(1.8), Inches(11), Inches(0.8), "Clone-Xs",
            size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(2.7), Inches(11), Inches(0.5),
            "Enterprise Unity Catalog Cloning — Made Simple",
            size=20, color=CRIMSON, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(3.5), Inches(11), Inches(0.4),
            "56 CLI Commands  •  31 UI Pages  •  61+ API Endpoints  •  88 Python Modules",
            size=14, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

# Links
add_rounded_rect(s, Inches(4.0), Inches(4.2), Inches(2.5), Inches(0.4), RGBColor(0x2A, 0x2A, 0x3E))
add_textbox(s, Inches(4.0), Inches(4.2), Inches(2.5), Inches(0.4),
            "github.com/viral0216/clone-xs", size=10, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)
add_rounded_rect(s, Inches(6.8), Inches(4.2), Inches(2.0), Inches(0.4), RGBColor(0x2A, 0x2A, 0x3E))
add_textbox(s, Inches(6.8), Inches(4.2), Inches(2.0), Inches(0.4),
            "MIT License", size=10, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

# Badge
add_rounded_rect(s, Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.4), CRIMSON)
add_textbox(s, Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.4),
            "Open Source — Start Cloning Today", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_slide_number(s, "14 / 14")

# ── Save ──
out_path = os.path.join(os.path.dirname(__file__), "Clone-Xs_Overview.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
