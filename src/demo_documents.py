"""Synthetic unstructured document generators for the /demo-data Documents tab.

Pairs with the existing structured demo generator (``src/demo_generator``)
and the streaming sibling (``src/demo_streaming``) — same opinionated
defaults, same Pydantic boundary, but a different output shape: the
generators here emit individual file bytes (PDF, Word, PowerPoint,
Excel, .eml) for RAG / GenAI demos that need a corpus of documents
rather than rows in a Delta table.

Architecture mirrors ``demo_streaming``:

    1. A registry (``DOCUMENT_TYPES``) maps operator-facing type IDs
       to ``(category_label, generator_fn)``.
    2. Each generator has the same signature
       ``(industry, faker, ai_client) -> tuple[bytes, dict]`` — the
       caller doesn't need to know which library produced the bytes.
    3. ``generate_documents`` is the top-level orchestrator that the
       API router invokes. It auto-creates the destination Volume,
       loops over (type, count) pairs, uploads bytes via
       ``client.files.upload``, and optionally creates / writes to a
       per-tab catalog (``volume_with_catalog``) or direct (``direct_table``)
       Delta table.

Why the dependencies are lazy-imported: the heavy doc deps
(``reportlab``, ``python-docx``, ``python-pptx``, ``openpyxl``) ship
in the ``[documents]`` optional extra. The router checks
:data:`DOCUMENTS_AVAILABLE` before dispatch and surfaces a clean 503
JSON response when the extra isn't installed — without crashing the
API server at import time.
"""

from __future__ import annotations

import io
import json
import logging
import random
import uuid
from collections.abc import Callable
from datetime import datetime, timedelta, timezone
from typing import Any

from databricks.sdk import WorkspaceClient

from src.client import execute_sql

logger = logging.getLogger(__name__)


# ── Lazy import probe ──────────────────────────────────────────────
#
# Each backend lib is probed once at module import. The probe lets the
# router decide whether to dispatch (deps present) or return a clean
# 503 (deps missing) without ever touching the generators.
DOCUMENTS_AVAILABLE: bool = False
_UNAVAILABLE_REASON: str | None = (
    "The `[documents]` optional extra is not installed. Run "
    "`pip install clone-xs[documents]` (which pulls reportlab, "
    "python-docx, python-pptx, openpyxl) and restart the API server."
)
try:
    import importlib.util as _ilutil

    if all(_ilutil.find_spec(m) is not None for m in ("reportlab", "docx", "pptx", "openpyxl")):
        DOCUMENTS_AVAILABLE = True
        _UNAVAILABLE_REASON = None
except Exception as _e:  # pragma: no cover — defensive
    _UNAVAILABLE_REASON = f"document deps probe failed: {_e}"


def is_available() -> tuple[bool, str | None]:
    """Return ``(available, reason)`` so the router can render an
    install hint when the extra isn't installed without itself
    crashing on the import."""
    return DOCUMENTS_AVAILABLE, _UNAVAILABLE_REASON


# ── AI drafting adapter ────────────────────────────────────────────
#
# Wraps the existing ``src.ai_service.AIService`` with a single
# ``draft(prompt, fallback)`` method that returns a one-or-two
# sentence narrative blurb. Routes through the user-picked
# Databricks Model Serving endpoint when one is available
# (``X-Databricks-Model`` header is the source of truth — clone-xs
# already sets this in ``ui/src/lib/api-client.ts`` from
# ``localStorage.dbx_model``); otherwise falls back to the Anthropic
# API path. The orchestrator constructs one of these per job and
# passes it into every generator.


# Delegate to the shared adapter so all five unstructured-tab
# generators (documents, media, knowledge, logs, code) go through the
# same backend selection, token budgeting, and telemetry.
from src.ai_drafter import maybe_ai as _maybe_ai  # noqa: E402,F401


def _rotate(*variants: str) -> str:
    """Pick one of N phrase variants. ``random.choice`` over a tuple,
    spelled out so the call sites read as 'rotate phrasing'."""
    return random.choice(variants)


def _maybe_section(prob: float) -> bool:
    """``random.random() < prob`` — a readable shorthand for optional
    section inclusion."""
    return random.random() < prob


# ── Doc-type → generator registry ──────────────────────────────────
#
# Keep the operator-facing IDs short and lowercase — they appear in
# the UI's checkbox grid and in API request bodies. The category
# label groups them visually in the UI ("PDF", "Word", "Excel", "Email").
#
# `gen_fn` is stored as a string for lazy lookup so this module can be
# imported without the doc deps installed (the registry itself uses
# only stdlib types). The orchestrator resolves the function name to
# the actual callable via `globals()` at dispatch time.

# Each registry entry has:
#   - category / extension / gen_fn (unchanged from D1)
#   - industries: dict[str, str] mapping industry → display label.
#                 The special key "*" is the default label used for any
#                 industry not explicitly listed. ABSENCE of both the
#                 industry key AND "*" means the type is HIDDEN for that
#                 industry (e.g. pdf_claim only shows for healthcare and
#                 insurance — financial users don't see "claim form" in
#                 their picker).
#
# The legacy "label" field is kept on each entry pointing at the "*"
# fallback so any code that still reads `entry["label"]` directly
# (older callers, tests) still gets a sensible string.
DOCUMENT_TYPES: dict[str, dict[str, Any]] = {
    # ── PDF: claim form (healthcare + insurance only) ──
    "pdf_claim": {
        "category": "PDF",
        "label": "Claim form",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_claim",
        "industries": {
            "healthcare": "Healthcare claim form",
            "insurance": "Insurance claim form",
        },
    },
    # ── PDF: invoice — every industry, with industry-specific naming ──
    "pdf_invoice": {
        "category": "PDF",
        "label": "Invoice",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_invoice",
        "industries": {
            "*": "Invoice",
            "healthcare": "Medical invoice",
            "retail": "Sales invoice",
            "telecom": "Service bill",
            "energy": "Utility bill",
            "education": "Tuition invoice",
            "real_estate": "Service invoice",
            "logistics": "Freight invoice",
            "insurance": "Premium invoice",
        },
    },
    # ── PDF: long-form legal contract — every industry ──
    "pdf_contract": {
        "category": "PDF",
        "label": "Legal contract",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_contract",
        "industries": {
            "*": "Legal contract",
            "healthcare": "Patient consent form",
            "financial": "Loan agreement",
            "retail": "Supplier agreement",
            "telecom": "Service agreement",
            "manufacturing": "Supplier agreement",
            "energy": "Power purchase agreement",
            "education": "Enrollment agreement",
            "real_estate": "Lease agreement",
            "logistics": "Carrier contract",
            "insurance": "Policy document",
        },
    },
    # ── DOCX: business letter — every industry ──
    "docx_letter": {
        "category": "Word",
        "label": "Business letter",
        "extension": "docx",
        "gen_fn": "_gen_docx_letter",
        "industries": {
            "*": "Business letter",
            "healthcare": "Patient letter",
            "financial": "Customer letter",
            "retail": "Customer service letter",
            "telecom": "Customer letter",
            "manufacturing": "Vendor letter",
            "energy": "Customer letter",
            "education": "Acceptance letter",
            "real_estate": "Tenant letter",
            "logistics": "Customer letter",
            "insurance": "Policyholder letter",
        },
    },
    # ── DOCX: quarterly report — every industry ──
    "docx_report": {
        "category": "Word",
        "label": "Quarterly report",
        "extension": "docx",
        "gen_fn": "_gen_docx_report",
        "industries": {
            "*": "Quarterly report",
            "healthcare": "Quarterly clinical report",
            "financial": "Quarterly financial report",
            "retail": "Quarterly sales report",
            "telecom": "Quarterly KPI report",
            "manufacturing": "Quarterly production report",
            "energy": "Quarterly grid report",
            "education": "Department annual report",
            "real_estate": "Quarterly property report",
            "logistics": "Quarterly fleet report",
            "insurance": "Quarterly underwriting report",
        },
    },
    # ── PPTX: deck — every industry ──
    "pptx_deck": {
        "category": "PowerPoint",
        "label": "Pitch deck",
        "extension": "pptx",
        "gen_fn": "_gen_pptx_deck",
        "industries": {
            "*": "Pitch deck",
            "healthcare": "Clinical pitch deck",
            "financial": "Investor pitch deck",
            "retail": "Merchandising deck",
            "telecom": "Network strategy deck",
            "manufacturing": "Operations deck",
            "energy": "Energy strategy deck",
            "education": "Course pitch deck",
            "real_estate": "Brokerage deck",
            "logistics": "Operations deck",
            "insurance": "Investor deck",
        },
    },
    # ── XLSX: budget — every industry ──
    "xlsx_budget": {
        "category": "Excel",
        "label": "Budget spreadsheet",
        "extension": "xlsx",
        "gen_fn": "_gen_xlsx_budget",
        "industries": {
            "*": "Budget spreadsheet",
            "healthcare": "Department budget",
            "financial": "Budget spreadsheet",
            "retail": "Store budget",
            "telecom": "Operations budget",
            "manufacturing": "Plant budget",
            "energy": "Operations budget",
            "education": "Department budget",
            "real_estate": "Property budget",
            "logistics": "Fleet budget",
            "insurance": "Operations budget",
        },
    },
    # ── XLSX: inventory — every industry ──
    "xlsx_inventory": {
        "category": "Excel",
        "label": "Inventory list",
        "extension": "xlsx",
        "gen_fn": "_gen_xlsx_inventory",
        "industries": {
            "*": "Inventory list",
            "healthcare": "Pharmacy inventory",
            "financial": "Asset register",
            "retail": "Store inventory",
            "telecom": "Network asset register",
            "manufacturing": "Parts inventory",
            "energy": "Asset register",
            "education": "Library catalog",
            "real_estate": "Property list",
            "logistics": "Fleet inventory",
            "insurance": "Asset register",
        },
    },
    # ── EML: message — every industry ──
    "eml_message": {
        "category": "Email",
        "label": "Email .eml",
        "extension": "eml",
        "gen_fn": "_gen_eml_message",
        "industries": {
            "*": "Email",
            "healthcare": "Care team email",
            "education": "Faculty email",
            "real_estate": "Agent email",
        },
    },
    # ── Industry-specific NEW types (one or two per industry).
    # All use shared layout helpers (_pdf_table_doc, _docx_letter_doc,
    # _xlsx_table_doc) defined below — keeps the per-type code small. ──
    # Healthcare
    "pdf_lab_report": {
        "category": "PDF",
        "label": "Lab report",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_lab_report",
        "industries": {"healthcare": "Lab report"},
    },
    "pdf_discharge_summary": {
        "category": "PDF",
        "label": "Patient discharge summary",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_discharge_summary",
        "industries": {"healthcare": "Patient discharge summary"},
    },
    # Financial
    "pdf_account_statement": {
        "category": "PDF",
        "label": "Account statement",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_account_statement",
        "industries": {"financial": "Account statement"},
    },
    "pdf_wire_confirmation": {
        "category": "PDF",
        "label": "Wire transfer confirmation",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_wire_confirmation",
        "industries": {"financial": "Wire transfer confirmation"},
    },
    # Retail
    "pdf_purchase_order": {
        "category": "PDF",
        "label": "Purchase order",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_purchase_order",
        "industries": {"retail": "Purchase order", "manufacturing": "Purchase order"},
    },
    "pdf_receipt": {
        "category": "PDF",
        "label": "Sales receipt",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_receipt",
        "industries": {"retail": "Sales receipt"},
    },
    # Telecom
    "pdf_sla_report": {
        "category": "PDF",
        "label": "SLA report",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_sla_report",
        "industries": {"telecom": "SLA report"},
    },
    "docx_outage_notice": {
        "category": "Word",
        "label": "Outage notification",
        "extension": "docx",
        "gen_fn": "_gen_docx_outage_notice",
        "industries": {"telecom": "Outage notification", "energy": "Outage notification"},
    },
    # Manufacturing
    "pdf_bom": {
        "category": "PDF",
        "label": "Bill of materials",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_bom",
        "industries": {"manufacturing": "Bill of materials"},
    },
    "pdf_qa_report": {
        "category": "PDF",
        "label": "Quality inspection report",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_qa_report",
        "industries": {"manufacturing": "Quality inspection report"},
    },
    # Energy
    "pdf_meter_reading": {
        "category": "PDF",
        "label": "Meter reading report",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_meter_reading",
        "industries": {"energy": "Meter reading report"},
    },
    # Education
    "pdf_transcript": {
        "category": "PDF",
        "label": "Academic transcript",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_transcript",
        "industries": {"education": "Academic transcript"},
    },
    "docx_syllabus": {
        "category": "Word",
        "label": "Course syllabus",
        "extension": "docx",
        "gen_fn": "_gen_docx_syllabus",
        "industries": {"education": "Course syllabus"},
    },
    # Real estate
    "pdf_property_listing": {
        "category": "PDF",
        "label": "Property listing",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_property_listing",
        "industries": {"real_estate": "Property listing"},
    },
    "pdf_disclosure": {
        "category": "PDF",
        "label": "Property disclosure",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_disclosure",
        "industries": {"real_estate": "Property disclosure"},
    },
    # Logistics
    "pdf_bol": {
        "category": "PDF",
        "label": "Bill of lading",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_bol",
        "industries": {"logistics": "Bill of lading"},
    },
    "pdf_customs": {
        "category": "PDF",
        "label": "Customs declaration",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_customs",
        "industries": {"logistics": "Customs declaration"},
    },
    # Insurance
    "pdf_underwriting_report": {
        "category": "PDF",
        "label": "Underwriting report",
        "extension": "pdf",
        "gen_fn": "_gen_pdf_underwriting_report",
        "industries": {"insurance": "Underwriting report"},
    },
    "docx_endorsement": {
        "category": "Word",
        "label": "Policy endorsement",
        "extension": "docx",
        "gen_fn": "_gen_docx_endorsement",
        "industries": {"insurance": "Policy endorsement"},
    },
}


# Average bytes per generated file, calibrated empirically by running
# each generator against `_AVG_BYTES_FIXTURE` and averaging the
# emitted size. Used by the preview endpoint so the UI can surface
# "approximately N MB total" before the operator clicks Generate.
# Stale-by-design: rough estimates beat exact sizes here because the
# UI just needs an order-of-magnitude.
_AVG_BYTES_PER_TYPE: dict[str, int] = {
    "pdf_claim": 3_500,
    "pdf_invoice": 2_800,
    "pdf_contract": 12_000,
    "docx_letter": 10_500,
    "docx_report": 18_000,
    "pptx_deck": 35_000,
    "xlsx_budget": 8_000,
    "xlsx_inventory": 9_000,
    "eml_message": 1_200,
    # New industry-specific types — calibrated against a smoke run.
    # Numbers are rough; preview tile only needs order-of-magnitude.
    "pdf_lab_report": 3_200,
    "pdf_discharge_summary": 6_500,
    "pdf_account_statement": 4_200,
    "pdf_wire_confirmation": 2_400,
    "pdf_purchase_order": 3_000,
    "pdf_receipt": 1_800,
    "pdf_sla_report": 5_000,
    "docx_outage_notice": 9_000,
    "pdf_bom": 4_500,
    "pdf_qa_report": 5_500,
    "pdf_meter_reading": 3_000,
    "pdf_transcript": 4_000,
    "docx_syllabus": 11_000,
    "pdf_property_listing": 3_500,
    "pdf_disclosure": 4_500,
    "pdf_bol": 3_000,
    "pdf_customs": 3_200,
    "pdf_underwriting_report": 6_000,
    "docx_endorsement": 10_500,
}

# Per-type throughput, calibrated on a developer machine. Used by the
# preview endpoint to estimate wall-time. PDF is the slowest (reportlab
# composes pages); XLSX is mid; .eml is essentially instant.
_GEN_PER_SECOND_PER_TYPE: dict[str, int] = {
    "pdf_claim": 200,
    "pdf_invoice": 220,
    "pdf_contract": 120,
    "docx_letter": 180,
    "docx_report": 150,
    "pptx_deck": 80,
    "xlsx_budget": 150,
    "xlsx_inventory": 130,
    "eml_message": 1000,
    # New types — most reuse the table-PDF helper, similar throughput
    # to pdf_invoice. Letter-style DOCX matches docx_letter's rate.
    "pdf_lab_report": 200,
    "pdf_discharge_summary": 160,
    "pdf_account_statement": 200,
    "pdf_wire_confirmation": 230,
    "pdf_purchase_order": 220,
    "pdf_receipt": 240,
    "pdf_sla_report": 180,
    "docx_outage_notice": 180,
    "pdf_bom": 200,
    "pdf_qa_report": 180,
    "pdf_meter_reading": 220,
    "pdf_transcript": 200,
    "docx_syllabus": 160,
    "pdf_property_listing": 220,
    "pdf_disclosure": 180,
    "pdf_bol": 220,
    "pdf_customs": 220,
    "pdf_underwriting_report": 160,
    "docx_endorsement": 170,
}


# ── Shared template content ───────────────────────────────────────
#
# Per-industry phrase pools driving the body text of generated
# documents. Kept in code (not YAML / template files) for v1 — the
# v2 roadmap moves these to a `src/demo_documents_templates/` dir
# so operators can extend without forking. Today the pools are small
# (~5 phrases per industry) which is enough for demo realism without
# needing Faker pools.

_INDUSTRY_CONTEXT: dict[str, dict[str, list[str]]] = {
    "healthcare": {
        "claim_diagnoses": [
            "I10 (Essential hypertension)",
            "E11.9 (Type 2 diabetes mellitus without complications)",
            "M54.5 (Low back pain)",
            "J45.909 (Asthma, unspecified)",
            "F33.1 (Major depressive disorder, recurrent, moderate)",
            "K21.9 (Gastro-esophageal reflux disease without esophagitis)",
            "N39.0 (Urinary tract infection, site not specified)",
            "M17.11 (Osteoarthritis, unilateral, primary)",
            "G47.33 (Obstructive sleep apnea (adult))",
            "E78.5 (Hyperlipidemia, unspecified)",
            "I25.10 (Atherosclerotic heart disease)",
            "F41.1 (Generalized anxiety disorder)",
            "M79.3 (Panniculitis, unspecified)",
            "L40.0 (Psoriasis vulgaris)",
            "H40.11X1 (Primary open-angle glaucoma, mild)",
            "K58.9 (Irritable bowel syndrome)",
            "N40.0 (Benign prostatic hyperplasia)",
            "R51 (Headache)",
        ],
        "treatment_codes": [
            "99213 — Office visit, established patient, low complexity",
            "97110 — Therapeutic exercises, 15 minutes",
            "85025 — Complete blood count with differential",
            "93000 — Electrocardiogram",
            "96372 — Therapeutic injection",
            "99214 — Office visit, established patient, moderate complexity",
            "20610 — Joint aspiration / injection, major joint",
            "73610 — X-ray ankle, complete",
            "99396 — Periodic preventive medicine, established patient",
            "90471 — Immunization administration, single",
            "94010 — Spirometry breathing test",
            "82607 — Vitamin B-12 level",
            "83036 — Hemoglobin A1c blood test",
            "76830 — Transvaginal ultrasound",
            "29881 — Arthroscopy, knee, surgical",
            "99203 — Office visit, new patient, moderate complexity",
            "27447 — Total knee arthroplasty",
            "92012 — Eye exam, established patient",
        ],
        "department_names": [
            "Internal Medicine",
            "Cardiology",
            "Orthopedics",
            "Family Medicine",
            "Pediatrics",
            "Dermatology",
            "Neurology",
            "Endocrinology",
            "Gastroenterology",
            "Pulmonology",
            "Rheumatology",
            "Oncology",
            "Psychiatry",
            "Radiology",
            "Emergency Medicine",
            "Obstetrics & Gynecology",
        ],
    },
    "financial": {
        "transaction_types": [
            "ACH credit",
            "Wire transfer",
            "Card-not-present purchase",
            "POS purchase",
            "Mobile deposit",
            "Direct deposit (payroll)",
            "Bill pay",
            "P2P transfer (Zelle / Venmo)",
            "ATM withdrawal",
            "Recurring subscription charge",
            "Foreign currency conversion",
            "Cashier's check",
            "ACH debit (auto-pay)",
            "Returned item / NSF",
            "Investment account transfer",
            "Loan disbursement",
            "Refund / reversal",
        ],
        "fee_categories": [
            "Maintenance fee",
            "Wire transfer fee",
            "Overdraft fee",
            "ATM withdrawal fee",
            "Foreign transaction fee",
            "Stop payment fee",
            "Returned item fee",
            "Inactive account fee",
            "Cash advance fee",
            "Paper statement fee",
            "Out-of-network ATM fee",
            "Excess withdrawal fee",
            "Replacement card fee",
            "Late payment fee",
            "Account closure fee",
        ],
        "department_names": [
            "Retail Banking",
            "Wealth Management",
            "Treasury Services",
            "Risk & Compliance",
            "Card Services",
            "Commercial Lending",
            "Mortgage Operations",
            "Fraud Investigations",
            "Anti-Money Laundering",
            "Capital Markets",
            "Trade Finance",
            "Corporate Banking",
            "Branch Operations",
            "Digital Channels",
            "Investment Banking",
        ],
    },
    "retail": {
        "product_categories": [
            "Apparel",
            "Home goods",
            "Electronics",
            "Grocery",
            "Pharmacy",
            "Beauty & personal care",
            "Toys & games",
            "Sporting goods",
            "Automotive accessories",
            "Office supplies",
            "Pet supplies",
            "Books & media",
            "Furniture",
            "Garden & outdoor",
            "Baby & kids",
            "Footwear",
            "Jewelry & accessories",
        ],
        "store_codes": [f"STR-{i:04d}" for i in range(1001, 1031)],
        "department_names": [
            "Merchandising",
            "Store Operations",
            "Supply Chain",
            "E-commerce",
            "Customer Care",
            "Loss Prevention",
            "Visual Merchandising",
            "Loyalty & CRM",
            "Pricing & Promotions",
            "Inventory Planning",
            "Returns & Fulfillment",
            "Vendor Management",
            "Buying & Sourcing",
            "Store Design",
            "Workforce Planning",
        ],
    },
    "manufacturing": {
        "part_categories": [
            "Hydraulic components",
            "Electrical assemblies",
            "Machined housings",
            "Sensor modules",
            "Fasteners",
            "Bearings & bushings",
            "Pneumatic actuators",
            "Gaskets & seals",
            "Wire harnesses",
            "Sheet metal brackets",
            "Pumps & motors",
            "Power transmission belts",
            "Cast iron weldments",
            "PCB assemblies",
            "Stamped steel components",
            "Welded sub-assemblies",
            "Plastic injection-molded parts",
        ],
        "department_names": [
            "Production Line A",
            "Production Line B",
            "Quality Assurance",
            "Maintenance",
            "Procurement",
            "Manufacturing Engineering",
            "Plant Safety",
            "Tooling & Fixtures",
            "Continuous Improvement",
            "Materials Handling",
            "Shipping & Receiving",
            "Plant Engineering",
            "Test & Inspection",
            "Final Assembly",
            "Process Control",
        ],
    },
    "energy": {
        "asset_types": [
            "Wind turbine",
            "Solar array",
            "Substation",
            "Transformer",
            "Smart meter",
            "Battery storage system",
            "Hydroelectric generator",
            "Combined-cycle gas turbine",
            "Distribution feeder",
            "Reclosing breaker",
            "Capacitor bank",
            "Voltage regulator",
            "Pole-mounted transformer",
            "EV fast-charging station",
            "Underground cable vault",
            "SCADA RTU",
        ],
        "department_names": [
            "Generation",
            "Transmission",
            "Distribution",
            "Customer Service",
            "Field Operations",
            "Grid Operations",
            "Renewables Integration",
            "Trading & Settlements",
            "Outage Management",
            "Asset Management",
            "Regulatory Affairs",
            "Vegetation Management",
            "Demand Response",
            "Metering Services",
            "System Planning",
        ],
        "outage_causes": [
            "Vegetation contact",
            "Equipment failure",
            "Severe weather",
            "Scheduled maintenance",
            "Third-party damage",
            "Animal contact (squirrel / bird)",
            "Vehicle collision with pole",
            "Lightning strike",
            "Substation transformer fault",
            "Underground cable failure",
            "Ice loading on conductors",
            "Wildfire-related precautionary shutoff",
            "Customer-side equipment fault",
            "Generation trip",
            "Cyber-physical event",
        ],
        "regulatory_bodies": [
            "FERC",
            "NERC",
            "PUC",
            "ISO/RTO",
            "Local utility commission",
            "EPA",
            "NRC",
            "DOE",
            "State energy commission",
            "Regional reliability council",
        ],
    },
    "telecom": {
        "service_types": [
            "Fiber broadband (1 Gbps)",
            "5G mobile postpaid",
            "Business VoIP",
            "MPLS circuit",
            "Cloud PBX",
            "SD-WAN service",
            "Dedicated Ethernet (10 Gbps)",
            "Wavelength service",
            "Hosted contact center",
            "Mobile prepaid",
            "Fixed wireless access",
            "Satellite backhaul",
            "Carrier Ethernet",
            "IoT connectivity (NB-IoT)",
            "Private 5G campus",
        ],
        "department_names": [
            "Network Operations",
            "Customer Care",
            "Field Engineering",
            "Wholesale & Carrier",
            "Product Management",
            "Service Provisioning",
            "Tower & Site Maintenance",
            "OSS / BSS Engineering",
            "Spectrum Management",
            "Capacity Planning",
            "Roaming & Interconnect",
            "Devices & Activations",
            "Billing Operations",
            "Network Security (SOC)",
        ],
        "sla_metrics": [
            "Network availability ≥ 99.95%",
            "MTTR ≤ 4 hours",
            "P1 incident response ≤ 15 minutes",
            "Latency ≤ 25 ms RTT",
            "Packet loss ≤ 0.1%",
            "Jitter ≤ 5 ms",
            "First-call resolution ≥ 80%",
            "Provisioning lead time ≤ 5 business days",
            "MTBF ≥ 8,760 hours",
            "Trouble-ticket close-out ≤ 24 hours",
            "Service credit threshold breach < 0.5% / month",
            "Dropped-call rate ≤ 0.5%",
            "5G coverage ≥ 92% of footprint",
        ],
        "outage_causes": [
            "Fiber cut",
            "BGP misconfiguration",
            "Power failure at POP",
            "DDoS attack",
            "Scheduled software upgrade",
            "Cell-site backhaul failure",
            "Tower antenna fault",
            "DNS resolver outage",
            "Authentication system failure",
            "Routing loop / black hole",
            "Capacity exhaustion (peak hour)",
            "Vendor / third-party dependency outage",
            "Battery backup depletion",
            "Cooling failure in equipment room",
        ],
    },
    "education": {
        "course_codes": [
            "CS101 — Introduction to Computer Science",
            "MATH240 — Linear Algebra",
            "ENG210 — Modern Literature",
            "BIO340 — Molecular Biology",
            "ECON101 — Microeconomics",
            "PHYS230 — Classical Mechanics",
            "CHEM150 — General Chemistry",
            "PSYC100 — Introduction to Psychology",
            "HIST220 — World History to 1500",
            "STAT310 — Statistical Inference",
            "PHIL120 — Introduction to Ethics",
            "ART180 — Studio Drawing",
            "POLS101 — American Government",
            "MUSC130 — Music Theory I",
            "GEOG260 — Physical Geography",
            "EDUC315 — Curriculum Design",
            "FREN201 — Intermediate French",
            "CS342 — Algorithms & Data Structures",
        ],
        "grade_letters": [
            "A",
            "A-",
            "B+",
            "B",
            "B-",
            "C+",
            "C",
            "C-",
            "D+",
            "D",
            "P",
            "NP",
            "I",
            "W",
            "AU",
        ],
        "department_names": [
            "Admissions",
            "Registrar",
            "Financial Aid",
            "Academic Affairs",
            "Student Services",
            "Bursar's Office",
            "Career Services",
            "International Student Office",
            "Library Services",
            "Athletics",
            "Information Technology",
            "Residential Life",
            "Counseling Center",
            "Honor Code Office",
            "Continuing Education",
        ],
        "academic_terms": [
            "Fall 2025",
            "Spring 2026",
            "Summer 2026",
            "Fall 2026",
            "Spring 2027",
            "Summer 2027",
            "Fall 2024",
            "Spring 2025",
            "Winter 2026",
            "Trimester 1 2026",
            "Trimester 2 2026",
            "Trimester 3 2026",
        ],
    },
    "real_estate": {
        "property_types": [
            "Single-family residence",
            "Condominium unit",
            "Multi-family duplex",
            "Commercial retail space",
            "Mixed-use development",
            "Townhouse",
            "Co-op apartment",
            "Industrial warehouse",
            "Office building (Class A)",
            "Office building (Class B)",
            "Vacant land — residential",
            "Vacant land — commercial",
            "Multi-family fourplex",
            "Manufactured home",
            "Hotel / hospitality property",
            "Self-storage facility",
        ],
        "department_names": [
            "Brokerage",
            "Property Management",
            "Asset Management",
            "Acquisitions",
            "Leasing",
            "Construction Management",
            "Title & Escrow",
            "Mortgage Origination",
            "Investor Relations",
            "Marketing",
            "Compliance",
            "Tenant Services",
            "Facilities & Maintenance",
            "Appraisal Services",
        ],
        "lease_terms": [
            "12-month standard residential lease",
            "24-month commercial lease with CPI escalation",
            "Month-to-month tenancy at will",
            "5-year triple-net commercial lease",
            "6-month corporate furnished lease",
            "10-year ground lease",
            "Short-term vacation rental (30 days)",
            "Office sublease through end of master term",
            "Modified gross lease, 3-year",
            "Percentage rent retail lease",
            "Build-to-suit lease, 15-year",
            "Co-working seat license, monthly",
            "Loft conversion, 18-month",
            "Single-tenant net lease, 7-year",
        ],
        "disclosure_items": [
            "Lead-based paint disclosure (pre-1978 construction)",
            "Radon test results",
            "Flood zone designation",
            "Known prior water intrusion",
            "Septic / well system status",
            "Active mold remediation history",
            "Roof age and last replacement date",
            "HVAC system age and service history",
            "Foundation cracks or settlement",
            "HOA assessment history",
            "Easements and rights of way",
            "Boundary survey discrepancies",
            "Past insurance claims",
            "Underground storage tank presence",
            "Historic-district designation",
            "Wetlands or environmental restrictions",
        ],
    },
    "logistics": {
        "freight_classes": [
            "Class 50 — clean freight",
            "Class 60 — bricks / cement",
            "Class 70 — car parts / engines",
            "Class 85 — crated machinery",
            "Class 100 — assembled goods",
            "Class 125 — small appliances",
            "Class 150 — auto sheet metal parts",
            "Class 175 — clothing / soft goods",
            "Class 200 — sheet metal parts, packaged TVs",
            "Class 250 — refrigerators / mattresses",
            "Class 300 — wood cabinets / tables",
            "Class 400 — deer antlers / oddly shaped",
            "Class 500 — bags of gold dust / ping-pong balls",
        ],
        "department_names": [
            "Dispatch",
            "Fleet Management",
            "Customs Brokerage",
            "Warehouse Operations",
            "Last-mile Delivery",
            "Yard Management",
            "Customer Service",
            "Linehaul Operations",
            "Driver Recruiting",
            "Safety & Compliance",
            "Pricing & Bid Management",
            "International Forwarding",
            "Cold Chain Operations",
            "Cross-dock Operations",
            "Reverse Logistics",
        ],
        "incoterms": [
            "EXW (Ex Works)",
            "FOB (Free on Board)",
            "CIF (Cost, Insurance, Freight)",
            "DDP (Delivered Duty Paid)",
            "FCA (Free Carrier)",
            "FAS (Free Alongside Ship)",
            "CFR (Cost and Freight)",
            "CPT (Carriage Paid To)",
            "CIP (Carriage and Insurance Paid)",
            "DAP (Delivered at Place)",
            "DPU (Delivered at Place Unloaded)",
        ],
        "damage_causes": [
            "Crushed during transit",
            "Water damage in container",
            "Improper securing",
            "Forklift puncture",
            "Temperature excursion",
            "Pallet collapse",
            "Top-load damage",
            "Concealed damage discovered at receiver",
            "Theft / pilferage",
            "Contamination from co-loaded freight",
            "Saltwater exposure (ocean freight)",
            "Vibration / road shock",
            "Improper labeling / mis-route",
            "Tipped on side during cornering",
        ],
    },
    "insurance": {
        "policy_types": [
            "Auto liability — comprehensive",
            "Homeowners HO-3",
            "Commercial general liability",
            "Term life — 20-year",
            "Workers' compensation",
            "Whole life",
            "Universal life",
            "Renters HO-4",
            "Condo HO-6",
            "Umbrella liability ($1M)",
            "Disability income",
            "Long-term care",
            "Commercial auto",
            "Cyber liability",
            "Directors & Officers (D&O)",
            "Errors & Omissions (E&O)",
            "Pet insurance",
            "Travel insurance",
        ],
        "claim_categories": [
            "Property damage — vehicle collision",
            "Bodily injury — third party",
            "Theft / burglary",
            "Natural disaster (hail / wind)",
            "Liability / litigation",
            "Fire damage",
            "Water damage (non-flood)",
            "Vandalism",
            "Hit and run — uninsured motorist",
            "Slip and fall — premises liability",
            "Wage loss / temporary disability",
            "Pet liability claim",
            "Cyber breach response",
            "Product liability",
            "Comprehensive theft of vehicle",
            "Hail damage to roof",
        ],
        "department_names": [
            "Underwriting",
            "Claims",
            "Actuarial",
            "Customer Service",
            "Compliance",
            "Special Investigations Unit (SIU)",
            "Reinsurance",
            "Subrogation",
            "Loss Control",
            "Agency Relations",
            "Catastrophe Response",
            "Medical Bill Review",
            "Litigation Management",
            "Product Development",
        ],
        "endorsement_codes": [
            "HO-15 — Special Form coverage extension",
            "PP-04-46 — Loss of Use enhancement",
            "CA-99-44 — Drive Other Car endorsement",
            "CGL-CG-2010 — Additional Insured",
            "WC-04-14 — Voluntary Compensation",
            "HO-04-90 — Personal Property Replacement Cost",
            "HO-04-95 — Water Backup of Sewer or Drain",
            "PP-03-08 — Auto Loan/Lease Coverage",
            "CGL-CG-2026 — Additional Insured (Designated Person)",
            "BP-04-17 — Hired Auto and Non-Owned Auto",
            "HO-04-42 — Permitted Incidental Occupancies",
            "CA-04-44 — Drive Other Car Broadened Coverage",
            "WC-00-04-14 — Notification of Change in Ownership",
            "CGL-CG-2503 — Designated Construction Project",
        ],
    },
    # Fallback for industries without explicit context — generators
    # use the structural defaults below.
    "_default": {
        "department_names": [
            "Operations",
            "Finance",
            "Sales",
            "Engineering",
            "Customer Success",
            "Human Resources",
            "Legal",
            "Marketing",
            "Information Technology",
            "Procurement",
            "Strategy",
            "Risk Management",
            "Internal Audit",
        ],
    },
}


def _ctx(industry: str, key: str, fallback: list[str] | None = None) -> list[str]:
    """Look up a context list for ``industry`` + ``key``. Falls back
    to the ``_default`` industry when the requested industry doesn't
    have the key, then to the supplied ``fallback`` list."""
    by_industry = _INDUSTRY_CONTEXT.get(industry, {})
    if key in by_industry:
        return by_industry[key]
    by_default = _INDUSTRY_CONTEXT["_default"]
    if key in by_default:
        return by_default[key]
    return fallback if fallback is not None else ["—"]


def label_for(type_id: str, industry: str) -> str:
    """Resolve the display label for a document type within an industry.

    Lookup order: the type's industry-specific label, then the type's
    ``"*"`` default, then the legacy ``label`` field, then the type id
    itself. Mirrors how :func:`types_for_industry` decides whether a
    type is shown — types without an explicit industry entry AND
    without ``"*"`` are hidden, but if a caller wants the label
    anyway (e.g. surfacing a stored count keyed by a hidden type) the
    ``label`` fallback gives a sensible string.
    """
    entry = DOCUMENT_TYPES.get(type_id) or {}
    industries: dict[str, str] = entry.get("industries") or {}
    if industry in industries:
        return industries[industry]
    if "*" in industries:
        return industries["*"]
    return str(entry.get("label") or type_id)


def types_for_industry(industry: str) -> list[dict[str, str]]:
    """Return the document-type list visible to ``industry``.

    Each entry is a flat dict with ``id``, ``category``, ``extension``,
    and ``label`` (already resolved for the industry). Used by the
    ``GET /demo-documents/types`` endpoint and by Pydantic validation
    to decide whether a count key is allowed for the chosen industry.

    A type is visible if its ``industries`` map contains an explicit
    entry for ``industry`` OR a ``"*"`` default entry. Order matches
    the registry insertion order (Python 3.7+ dict ordering), so
    callers get a stable picker layout.
    """
    out: list[dict[str, str]] = []
    for type_id, entry in DOCUMENT_TYPES.items():
        industries: dict[str, str] = entry.get("industries") or {}
        if industry not in industries and "*" not in industries:
            continue
        out.append(
            {
                "id": type_id,
                "category": str(entry["category"]),
                "extension": str(entry["extension"]),
                "label": label_for(type_id, industry),
            }
        )
    return out


# ── Per-type generators ───────────────────────────────────────────
#
# Each generator returns ``(file_bytes, metadata_dict)``. The metadata
# is what gets serialized into the catalog table's `metadata_json`
# column AND surfaces in the UI's per-cell summary. Keep keys
# JSON-serializable (no datetime objects, no numpy types).


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def _gen_pdf_claim(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a healthcare-style claim form as PDF bytes.

    Layout varies per call:
      * 1-3 procedure rows (random sample of CPT codes)
      * Optional "Notes from provider" narrative (40% of the time, AI-drafted when on)
      * Optional "Patient consent" block (30% of the time)
      * Heading variant (3 distinct title rotations)
      * Footer disclaimer rotated across 4 variants
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        Table,
        TableStyle,
    )
    from reportlab.lib import colors

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, title="Claim form")
    styles = getSampleStyleSheet()
    story: list[Any] = []

    claim_id = f"CLM-{uuid.uuid4().hex[:10].upper()}"
    patient_name = fkr.name()
    provider_name = f"Dr. {fkr.last_name()}"
    facility_suffix = _rotate(
        "Medical Group", "Health Partners", "Clinic Associates", "Specialty Care", "Family Practice"
    )
    facility = fkr.company() + " " + facility_suffix
    date_of_service = fkr.date_between(start_date="-180d", end_date="today").isoformat()
    diagnosis = random.choice(_ctx(industry, "claim_diagnoses", ["—"]))

    title_variant = _rotate(
        "Claim Form (Demo)",
        "Health Insurance Claim",
        "Medical Services Claim",
        "Provider Claim Submission",
    )

    story.append(Paragraph(facility, styles["Title"]))
    story.append(Paragraph(title_variant, styles["Heading2"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"<b>Claim ID:</b> {claim_id}", styles["Normal"]))
    story.append(Paragraph(f"<b>Patient:</b> {patient_name}", styles["Normal"]))
    story.append(Paragraph(f"<b>Provider:</b> {provider_name}", styles["Normal"]))
    story.append(Paragraph(f"<b>Date of Service:</b> {date_of_service}", styles["Normal"]))
    story.append(Paragraph(f"<b>Primary Diagnosis:</b> {diagnosis}", styles["Normal"]))
    story.append(Spacer(1, 18))

    story.append(Paragraph("<b>Procedures</b>", styles["Heading3"]))
    pool = _ctx(industry, "treatment_codes", ["—"])
    treatments = random.sample(pool, k=min(random.randint(1, 4), len(pool)))
    rows = [["CPT / Code", "Description", "Charge"]]
    total = 0
    for t in treatments:
        amount = round(random.uniform(75, 850), 2)
        total += amount
        # Split "code — description" if present
        if " — " in t:
            code, desc = t.split(" — ", 1)
        else:
            code, desc = t, ""
        rows.append([code, desc, f"${amount:.2f}"])
    rows.append(["", "Total", f"${total:.2f}"])

    table = Table(rows, colWidths=[100, 280, 80])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                ("FONTNAME", (0, -1), (-1, -1), "Helvetica-Bold"),
                ("ALIGN", (-1, 1), (-1, -1), "RIGHT"),
            ]
        )
    )
    story.append(table)

    # Optional "Notes from provider" — AI-drafted narrative when AI mode on.
    if _maybe_section(0.4):
        story.append(Spacer(1, 14))
        story.append(Paragraph("<b>Notes from provider</b>", styles["Heading3"]))
        notes = _maybe_ai(
            ai_client,
            (
                f"Write 2 short clinical sentences from {provider_name} explaining the "
                f"medical necessity of the procedures for diagnosis '{diagnosis}'. "
                f"Patient: {patient_name}. Date of service: {date_of_service}. "
                f"Tone: clinical, concise, third person."
            ),
            fallback=_rotate(
                f"Patient was evaluated for {diagnosis.split('(')[0].strip()} on {date_of_service}; treatment plan documented in chart.",
                f"Continued management of {diagnosis.split('(')[0].strip()}. Procedures performed in accordance with current standard of care.",
                "Follow-up visit; clinical findings consistent with prior assessment. Patient tolerated procedures well.",
                "Initial workup completed. Will reassess at next scheduled visit.",
            ),
            max_tokens=120,
        )
        story.append(Paragraph(notes, styles["Normal"]))

    # Optional "Patient consent" — 30% of the time.
    if _maybe_section(0.3):
        story.append(Spacer(1, 12))
        story.append(Paragraph("<b>Patient consent</b>", styles["Heading3"]))
        story.append(
            Paragraph(
                _rotate(
                    "Patient acknowledges receipt of treatment and authorises release of records to insurer.",
                    "Signed informed-consent form on file; copy attached to chart.",
                    "Patient has been informed of treatment plan and associated risks; verbal consent obtained.",
                ),
                styles["Normal"],
            )
        )

    story.append(Spacer(1, 24))
    story.append(
        Paragraph(
            "<i>"
            + _rotate(
                "Generated by Clone-Xs demo data — not a real claim. All names, IDs, and amounts are synthetic.",
                "DEMO ONLY — synthetic data. Not for clinical use, billing, or patient care.",
                "This is a fictitious claim form generated for demonstration purposes. No real patient information is contained herein.",
                "Sample claim form — not for production use. All fields populated from synthetic data sources.",
            )
            + "</i>",
            styles["Italic"],
        )
    )

    doc.build(story)
    return buf.getvalue(), {
        "claim_id": claim_id,
        "patient_name": patient_name,
        "provider_name": provider_name,
        "diagnosis": diagnosis,
        "total_charges": total,
        "procedure_count": len(treatments),
        "page_count": 1,  # SimpleDocTemplate auto-paginates; 1 in practice for this content
    }


def _gen_pdf_invoice(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a B2B invoice as PDF bytes.

    Variations: 2-8 line items, optional payment-terms paragraph
    (60% of the time, AI-drafted when on), rotating tax-rate footer,
    optional purchase-order reference line.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib import colors

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, title="Invoice")
    styles = getSampleStyleSheet()
    story: list[Any] = []

    invoice_id = f"INV-{datetime.now().strftime('%Y%m')}-{random.randint(10000, 99999)}"
    vendor = fkr.company()
    customer = fkr.company()
    issue_date = fkr.date_between(start_date="-60d", end_date="today").isoformat()
    due_date = fkr.date_between(start_date="today", end_date="+30d").isoformat()
    po_ref = f"PO-{random.randint(100000, 999999)}" if _maybe_section(0.5) else None

    story.append(Paragraph(vendor, styles["Title"]))
    story.append(
        Paragraph(
            _rotate(
                f"INVOICE — {invoice_id}",
                f"BILL — {invoice_id}",
                f"STATEMENT OF CHARGES — {invoice_id}",
            ),
            styles["Heading2"],
        )
    )
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"<b>Bill To:</b> {customer}", styles["Normal"]))
    story.append(Paragraph(f"<b>Issue Date:</b> {issue_date}", styles["Normal"]))
    story.append(Paragraph(f"<b>Due Date:</b> {due_date}", styles["Normal"]))
    if po_ref:
        story.append(Paragraph(f"<b>Customer PO:</b> {po_ref}", styles["Normal"]))
    story.append(Spacer(1, 18))

    rows = [["Item", "Qty", "Unit Price", "Subtotal"]]
    total = 0.0
    line_items = random.randint(2, 8)
    for _ in range(line_items):
        item = fkr.bs().capitalize()
        qty = random.randint(1, 25)
        unit_price = round(random.uniform(20, 1200), 2)
        subtotal = round(qty * unit_price, 2)
        total += subtotal
        rows.append([item, str(qty), f"${unit_price:.2f}", f"${subtotal:.2f}"])
    tax_rate = random.choice([0.0625, 0.0725, 0.0825, 0.0975, 0.10])
    tax = round(total * tax_rate, 2)
    grand_total = round(total + tax, 2)
    rows.append(["", "", "Subtotal", f"${total:.2f}"])
    rows.append(["", "", f"Tax ({tax_rate * 100:.2f}%)", f"${tax:.2f}"])
    rows.append(["", "", "Total Due", f"${grand_total:.2f}"])

    table = Table(rows, colWidths=[260, 50, 90, 90])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("GRID", (0, 0), (-1, line_items), 0.25, colors.grey),
                ("FONTNAME", (-2, -1), (-1, -1), "Helvetica-Bold"),
                ("ALIGN", (-2, 1), (-1, -1), "RIGHT"),
            ]
        )
    )
    story.append(table)

    # Optional payment-terms / remittance paragraph.
    if _maybe_section(0.6):
        story.append(Spacer(1, 18))
        story.append(Paragraph("<b>Payment terms</b>", styles["Heading3"]))
        terms = _maybe_ai(
            ai_client,
            (
                f"Write 1-2 sentences of payment-terms language for invoice {invoice_id}, "
                f"due {due_date}, total ${grand_total}. Mention preferred payment method "
                f"(ACH or wire) and late-fee policy. Tone: professional, brief."
            ),
            fallback=_rotate(
                "Net-30 terms apply. ACH preferred; wire instructions available on request. Late payments accrue interest at 1.5% per month.",
                "Payment due within 30 days of issue. Remit by ACH transfer to the account on file. A 1.5% monthly service charge applies to past-due balances.",
                "Terms: Net-30. Payment may be made by ACH, wire, or check. Invoices unpaid after the due date are subject to a 1.5% monthly finance charge.",
            ),
            max_tokens=120,
        )
        story.append(Paragraph(terms, styles["Normal"]))

    doc.build(story)
    return buf.getvalue(), {
        "invoice_id": invoice_id,
        "vendor": vendor,
        "customer": customer,
        "line_items": line_items,
        "tax_rate": tax_rate,
        "po_ref": po_ref,
        "grand_total": grand_total,
        "page_count": 1,
    }


def _gen_pdf_contract(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a multi-page legal contract as PDF bytes.

    Templated structure: parties, recitals, terms (5–10 numbered
    sections), signature blocks. Roughly 3 pages of content, which
    is more useful than a one-pager for embeddings / chunking demos.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, title="Contract")
    styles = getSampleStyleSheet()
    story: list[Any] = []

    contract_id = f"CON-{uuid.uuid4().hex[:8].upper()}"
    party_a = fkr.company()
    party_b = fkr.company()
    effective_date = fkr.date_between(start_date="-365d", end_date="today").isoformat()

    title_variant = _rotate(
        "SERVICES AGREEMENT",
        "MASTER SERVICES AGREEMENT",
        "PROFESSIONAL SERVICES AGREEMENT",
        "CONSULTING SERVICES AGREEMENT",
    )
    story.append(Paragraph(title_variant, styles["Title"]))
    story.append(Spacer(1, 18))
    state = random.choice(
        ["Delaware", "California", "New York", "Texas", "Massachusetts", "Illinois", "Washington"]
    )
    venue_city = {
        "Delaware": "Wilmington",
        "California": "San Francisco",
        "New York": "New York",
        "Texas": "Austin",
        "Massachusetts": "Boston",
        "Illinois": "Chicago",
        "Washington": "Seattle",
    }[state]
    story.append(
        Paragraph(
            f"This {title_variant.title()} (the &quot;Agreement&quot;), dated as of "
            f"{effective_date} (the &quot;Effective Date&quot;), is entered "
            f"into by and between <b>{party_a}</b>, a corporation "
            f"(&quot;{party_a}&quot;), and <b>{party_b}</b>, a corporation "
            f"(&quot;{party_b}&quot;).",
            styles["Normal"],
        )
    )
    story.append(Spacer(1, 18))

    # All available sections; we'll pick a random subset for each contract.
    candidate_sections: list[tuple[str, str]] = [
        (
            "Scope of Services",
            _maybe_ai(
                ai_client,
                f"Write a 'Scope of Services' contract clause (2-3 sentences) for a {industry}-industry services agreement between {party_a} and {party_b}. Use formal contract language.",
                fallback="Provider shall perform the services described in Schedule A attached hereto. Such services shall be performed in a professional manner consistent with industry standards and applicable law.",
                max_tokens=140,
            ),
        ),
        (
            "Term",
            _maybe_ai(
                ai_client,
                f"Write a 'Term' contract clause (2 sentences) for the agreement effective {effective_date}. Include initial term length and renewal mechanism.",
                fallback="This Agreement shall commence on the Effective Date and continue for an initial term of twelve (12) months, automatically renewing for successive twelve-month periods unless terminated by either party with thirty (30) days written notice.",
                max_tokens=120,
            ),
        ),
        (
            "Compensation",
            "Customer shall pay Provider the fees set forth in Schedule B. Invoices shall be issued monthly and are due within thirty (30) days of receipt. Late payments shall accrue interest at 1.5% per month.",
        ),
        (
            "Confidentiality",
            _maybe_ai(
                ai_client,
                "Write a 'Confidentiality' contract clause (2-3 sentences) for a services agreement. Include survival period after termination.",
                fallback="Each party agrees to hold the other's confidential information in strict confidence and to use it solely for the purpose of performing this Agreement. This obligation shall survive termination for a period of three (3) years.",
                max_tokens=120,
            ),
        ),
        (
            "Limitation of Liability",
            "In no event shall either party be liable for indirect, special, or consequential damages, regardless of the form of action. Total liability shall not exceed the fees paid in the twelve months preceding the claim.",
        ),
        (
            "Governing Law",
            f"This Agreement shall be governed by the laws of the State of {state}, without regard to its conflict-of-laws principles. Any dispute shall be resolved exclusively in the state or federal courts located in {venue_city}, {state}.",
        ),
        (
            "Entire Agreement",
            "This Agreement constitutes the entire understanding between the parties and supersedes all prior negotiations, representations, and agreements, whether written or oral, with respect to the subject matter hereof.",
        ),
        (
            "Indemnification",
            "Each party shall indemnify and hold the other harmless from and against any third-party claims arising from the indemnifying party's gross negligence or wilful misconduct in the performance of this Agreement.",
        ),
        (
            "Insurance",
            "Provider shall maintain commercial general liability insurance with limits of not less than $1,000,000 per occurrence and $2,000,000 in the aggregate, naming Customer as an additional insured.",
        ),
        (
            "Force Majeure",
            "Neither party shall be liable for any failure or delay in performance under this Agreement to the extent such failure or delay is caused by an event beyond the reasonable control of such party, including acts of God, war, terrorism, fire, flood, or governmental action.",
        ),
        (
            "Independent Contractor",
            "The relationship of the parties shall be that of independent contractors. Nothing in this Agreement shall be construed to create a partnership, joint venture, agency, or employment relationship between the parties.",
        ),
        (
            "Assignment",
            "Neither party may assign this Agreement, in whole or in part, without the prior written consent of the other party, except that either party may assign this Agreement to a successor in connection with a merger, acquisition, or sale of substantially all of its assets.",
        ),
    ]
    # Always include the first 5 anchors; randomize the remaining picks.
    must_have = candidate_sections[:5]
    optional = candidate_sections[5:]
    extra_count = random.randint(2, len(optional))
    extras = random.sample(optional, k=extra_count)
    sections = must_have + extras
    random.shuffle(extras)  # shuffle just the extras so anchor order stays predictable

    for i, (title, body) in enumerate(sections, start=1):
        story.append(Paragraph(f"<b>{i}. {title}</b>", styles["Heading3"]))
        story.append(Paragraph(body, styles["Normal"]))
        story.append(Spacer(1, 12))

    story.append(PageBreak())
    story.append(Paragraph("IN WITNESS WHEREOF", styles["Heading3"]))
    story.append(
        Paragraph(
            "the parties have executed this Agreement as of the Effective Date.",
            styles["Normal"],
        )
    )
    story.append(Spacer(1, 30))
    for party in (party_a, party_b):
        story.append(Paragraph(f"<b>{party}</b>", styles["Normal"]))
        story.append(Paragraph("By: ___________________________", styles["Normal"]))
        story.append(Paragraph(f"Name: {fkr.name()}", styles["Normal"]))
        story.append(Paragraph(f"Title: {fkr.job()}", styles["Normal"]))
        story.append(Spacer(1, 24))

    doc.build(story)
    return buf.getvalue(), {
        "contract_id": contract_id,
        "party_a": party_a,
        "party_b": party_b,
        "effective_date": effective_date,
        "section_count": len(sections),
        "governing_state": state,
        "page_count": 2,
    }


def _gen_docx_letter(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a one-page business letter as DOCX bytes.

    Variations: 2-4 body paragraphs (each AI-drafted on demand),
    rotating sign-offs and salutations, optional CC line.
    """
    from docx import Document

    doc = Document()
    sender = fkr.name()
    sender_title = fkr.job()
    sender_company = fkr.company()
    recipient = fkr.name()
    recipient_company = fkr.company()
    today = datetime.now().strftime("%B %d, %Y")
    department = random.choice(_ctx(industry, "department_names", ["Operations"]))
    salutation = _rotate("Dear", "Hello", "Greetings,", "To whom it may concern:")
    closing = _rotate(
        "Sincerely,", "Best regards,", "Kind regards,", "Yours sincerely,", "Warm regards,"
    )
    body_paragraphs = random.randint(2, 4)
    body_topic_pool = [
        ("the upcoming engagement", "next-quarter planning"),
        ("recent contract renegotiation", "fee schedule alignment"),
        ("the integration kickoff", "stakeholder communications"),
        ("the audit findings", "remediation timelines"),
        ("the joint initiative", "executive sponsorship"),
    ]
    topic, sub_topic = random.choice(body_topic_pool)

    doc.add_paragraph(sender_company)
    doc.add_paragraph(today)
    doc.add_paragraph()
    doc.add_paragraph(f"{recipient}\n{recipient_company}")
    doc.add_paragraph()
    if salutation.endswith(":"):
        doc.add_paragraph(salutation)
    else:
        doc.add_paragraph(f"{salutation} {recipient.split()[0]},")

    # First body paragraph — context-setter.
    doc.add_paragraph(
        _maybe_ai(
            ai_client,
            f"Write the opening paragraph (2-3 sentences) of a formal business letter from {sender} ({sender_title}) at {sender_company} to {recipient} at {recipient_company}, regarding {topic} in the {department} group. Tone: professional and warm.",
            fallback=(
                f"I am writing to follow up on our recent discussion regarding {topic} between our organizations. "
                f"As you know, the next phase of work will require coordination across several teams, and I want "
                f"to ensure we are aligned on the timeline and deliverables for the {department.lower()} engagement."
            ),
            max_tokens=160,
        )
    )

    if body_paragraphs >= 2:
        doc.add_paragraph(
            _maybe_ai(
                ai_client,
                f"Write the middle paragraph (2 sentences) of the letter — discussing {sub_topic}. Tone: business-professional, no fluff.",
                fallback=(
                    "Per our conversation, we will share the next deliverable by the end of the quarter. "
                    "My team is reviewing the proposed approach and will have detailed feedback ready for our follow-up meeting."
                ),
                max_tokens=130,
            )
        )

    if body_paragraphs >= 3:
        doc.add_paragraph(
            _maybe_ai(
                ai_client,
                f"Write a brief paragraph (1-2 sentences) covering open action items related to {sub_topic}. Tone: collaborative.",
                fallback=(
                    "In parallel, we are reviewing the supporting documentation and will circulate any open questions ahead of next week's session."
                ),
                max_tokens=110,
            )
        )

    if body_paragraphs >= 4:
        doc.add_paragraph(
            _maybe_ai(
                ai_client,
                "Write a brief 'next steps' paragraph (1 sentence) wrapping up the letter.",
                fallback=(
                    "I will plan to circulate a summary of next steps following our discussion next week."
                ),
                max_tokens=90,
            )
        )

    # Closing paragraph — always.
    doc.add_paragraph(
        _rotate(
            "Please let me know if there is anything else I can do to move this forward.",
            "I appreciate your continued partnership and look forward to next steps.",
            "Thank you for your attention to this matter; I'm happy to discuss further at your convenience.",
            "Please feel free to reach out with any questions or clarifications.",
        )
    )
    doc.add_paragraph()
    doc.add_paragraph(closing)
    doc.add_paragraph()
    doc.add_paragraph(f"{sender}\n{sender_title}\n{sender_company}")

    cc_recipients: list[str] = []
    if _maybe_section(0.35):
        cc_recipients = [fkr.name() for _ in range(random.randint(1, 3))]
        doc.add_paragraph()
        doc.add_paragraph(f"cc: {', '.join(cc_recipients)}")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue(), {
        "sender": sender,
        "recipient": recipient,
        "department": department,
        "body_paragraphs": body_paragraphs,
        "topic": topic,
        "cc_count": len(cc_recipients),
        "page_count": 1,
    }


def _gen_docx_report(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a multi-page quarterly report as DOCX bytes."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    quarter = f"Q{random.randint(1, 4)} {random.randint(2023, 2025)}"
    department = random.choice(_ctx(industry, "department_names", ["Operations"]))
    author = fkr.name()
    company = fkr.company()

    doc.add_heading(f"{department} Quarterly Report — {quarter}", level=0)
    doc.add_paragraph(f"Prepared by {author}, {company}")
    doc.add_paragraph()

    revenue = random.randint(8, 25)
    yoy_dir = random.choice(["up", "down"])
    yoy_pct = random.randint(2, 18)
    margin = random.randint(15, 35)
    capex = random.randint(1, 8)
    bus = random.randint(3, 12)
    savings = random.randint(500, 4000)
    hires = random.randint(8, 25)

    candidate_sections: list[tuple[str, str]] = [
        (
            "Executive Summary",
            _maybe_ai(
                ai_client,
                f"Write a 3-4 sentence executive summary for the {department} group's {quarter} quarterly report. Mention revenue ${revenue}M ({yoy_dir} {yoy_pct}% YoY), margin {margin}%, key wins. Tone: confident but candid.",
                fallback=(
                    f"This report summarizes the {department} team's performance for {quarter}. "
                    f"Key initiatives delivered on time include improved customer satisfaction "
                    f"scores, on-target revenue milestones, and continued investment in operational "
                    f"efficiency. The team navigated several headwinds, particularly around supply "
                    f"chain volatility and shifting customer demand patterns."
                ),
                max_tokens=200,
            ),
        ),
        (
            "Financial Performance",
            f"Revenue for the quarter came in at ${revenue}M, "
            f"{yoy_dir} {yoy_pct}% year-over-year. "
            f"Operating margin was {margin}%, in line with the prior "
            f"quarter. Capital expenditure totaled ${capex}M, weighted "
            f"toward modernization of legacy infrastructure.",
        ),
        (
            "Key Initiatives",
            _maybe_ai(
                ai_client,
                f"Write a 3-sentence 'Key Initiatives' section for {department} in {quarter}. Reference: customer-data platform reached {bus} business units, cost-optimization saved ${savings}K, {hires} new hires.",
                fallback=(
                    f"Three multi-quarter initiatives advanced in {quarter}. The customer-data "
                    f"platform completed its initial rollout to {bus} business units. The "
                    f"cost-optimization program identified ${savings}K in annualized savings. "
                    f"The talent-development workstream onboarded {hires} new hires across "
                    f"engineering, product, and operations."
                ),
                max_tokens=180,
            ),
        ),
        (
            "Risk and Mitigation",
            _maybe_ai(
                ai_client,
                f"Write 3 sentences on risk areas the {department} team is monitoring this quarter. Cover: vendor concentration, regulatory shifts, turnover.",
                fallback=(
                    f"The team is monitoring three primary risk areas: vendor concentration in "
                    f"the supply chain, regulatory shifts impacting the {department.lower()} "
                    f"business, and turnover in critical roles. Mitigation plans are in place "
                    f"for each, with quarterly reviews at the leadership level."
                ),
                max_tokens=160,
            ),
        ),
        (
            "Outlook",
            _maybe_ai(
                ai_client,
                f"Write a 2-3 sentence 'Outlook' section for the next quarter for the {department} group. Mention strategic initiatives, platform expansion, expected revenue/margin direction.",
                fallback=(
                    "For the next quarter, the team will prioritize completion of the in-flight "
                    "strategic initiatives, expansion of the customer-data platform to remaining "
                    "business units, and continued investment in operational excellence. "
                    "We expect revenue and margin to remain in the current range."
                ),
                max_tokens=160,
            ),
        ),
        (
            "Customer Highlights",
            _maybe_ai(
                ai_client,
                f"Write a 2-sentence 'Customer Highlights' section for the {department} {quarter} report — name 2 plausible customer wins or expansions.",
                fallback="Two enterprise customers extended their contracts this quarter, including a multi-year renewal in the Western region. Net Promoter Score improved by three points quarter-over-quarter.",
                max_tokens=140,
            ),
        ),
        (
            "Operational Metrics",
            f"Service availability was {round(random.uniform(99.5, 99.99), 2)}%, exceeding the {round(random.uniform(99.5, 99.9), 2)}% commitment. Mean time to resolution improved {random.randint(5, 30)}% versus the prior quarter. Backlog closed at {random.randint(20, 200)} open items.",
        ),
        (
            "People & Talent",
            f"Headcount ended the quarter at {random.randint(40, 250)}, a net change of {random.randint(-5, 30)} from the prior quarter. Voluntary attrition was {round(random.uniform(2.0, 12.0), 1)}%, in line with industry benchmarks.",
        ),
    ]
    # Always Exec Summary + Outlook; pick 2-5 of the rest randomly.
    must_have_indices = [0, 4]  # Exec Summary, Outlook
    optional_indices = [i for i in range(len(candidate_sections)) if i not in must_have_indices]
    extra_count = random.randint(2, len(optional_indices))
    chosen_extra = sorted(random.sample(optional_indices, k=extra_count))
    chosen_indices = sorted(
        set(must_have_indices + chosen_extra), key=lambda i: (i != 0, i == 4, i)
    )
    sections = [candidate_sections[i] for i in chosen_indices]

    for title_text, body in sections:
        doc.add_heading(title_text, level=1)
        para = doc.add_paragraph(body)
        for run in para.runs:
            run.font.size = Pt(11)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue(), {
        "quarter": quarter,
        "department": department,
        "author": author,
        "revenue_m": revenue,
        "yoy_pct": yoy_pct,
        "yoy_dir": yoy_dir,
        "section_count": len(sections),
        "page_count": 3,
    }


def _gen_pptx_deck(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a 6-slide pitch deck as PPTX bytes."""
    from pptx import Presentation
    from pptx.util import Pt

    pres = Presentation()
    company = fkr.company()
    industry_label = industry.replace("_", " ").title()
    audience = fkr.company()

    # Title slide
    slide_layout = pres.slide_layouts[0]
    s = pres.slides.add_slide(slide_layout)
    s.shapes.title.text = f"{company}"
    s.placeholders[1].text = f"Customer briefing for {audience} — {industry_label}"

    # Bullets per content slide. We sample which slides are included
    # and the bullets within each so two decks for the same audience
    # never look identical.
    candidate_slides: list[tuple[str, list[str]]] = [
        (
            "The opportunity",
            [
                _maybe_ai(
                    ai_client,
                    f"Write a 1-line punchy bullet about structural headwinds in the {industry_label} sector for a sales deck.",
                    fallback=f"{industry_label} customers face structural headwinds in their data estates",
                    max_tokens=40,
                ),
                _maybe_ai(
                    ai_client,
                    f"Write a 1-line bullet about how legacy systems block modern {industry_label} data needs.",
                    fallback="Legacy systems can't support modern data needs",
                    max_tokens=40,
                ),
                f"${random.randint(40, 800)}M TAM in the next {random.choice([12, 18, 24, 36])} months",
            ],
        ),
        (
            "What we deliver",
            [
                f"Unified data platform across the {industry_label} estate",
                f"Real-time analytics with sub-{random.choice(['100ms', '500ms', 'second'])} latency",
                "Built-in AI capabilities for next-gen use cases",
            ],
        ),
        (
            "Customer outcomes",
            [
                f"~{random.randint(20, 45)}% reduction in time to insight",
                f"~{random.randint(30, 60)}% lower total cost of ownership",
                f"{random.choice(['Faster', 'Streamlined', 'Improved'])} regulatory reporting cycles",
            ],
        ),
        (
            "Why now",
            [
                "Modern lakehouse architecture is production-ready",
                "Generative AI is unlocking new use cases",
                _rotate(
                    "Regulatory pressure is accelerating modernization",
                    "Customer expectations have outpaced legacy capabilities",
                    "Cloud economics now favour consolidation over expansion",
                ),
            ],
        ),
        (
            "Next steps",
            [
                "Schedule technical deep-dive",
                f"Run a {random.choice([30, 60, 90])}-day proof of concept",
                "Identify executive sponsor for the engagement",
            ],
        ),
        (
            "Reference customers",
            [
                f"{fkr.company()} — {random.choice([5, 10, 15, 25])}x query speedup",
                f"{fkr.company()} — single source of truth across {random.randint(3, 12)} business units",
                f"{fkr.company()} — {round(random.uniform(1.5, 9.0), 1)}x cost reduction",
            ],
        ),
        (
            "Architecture at a glance",
            [
                "Bronze / silver / gold medallion layers in Unity Catalog",
                _rotate(
                    "Streaming + batch in one pipeline",
                    "Continuous ingestion via Auto Loader",
                    "Decoupled storage + compute",
                ),
                "Built-in lineage, governance, and PII detection",
            ],
        ),
    ]
    must_have = [
        candidate_slides[0],
        candidate_slides[1],
        candidate_slides[4],
    ]  # Opportunity, Deliver, Next steps
    optional = [s for i, s in enumerate(candidate_slides) if i not in (0, 1, 4)]
    extras = random.sample(optional, k=random.randint(2, len(optional)))
    bullets_per_slide = must_have[:2] + extras + must_have[2:]

    for title, bullets in bullets_per_slide:
        slide_layout = pres.slide_layouts[1]
        s = pres.slides.add_slide(slide_layout)
        s.shapes.title.text = title
        body = s.placeholders[1].text_frame
        body.text = bullets[0]
        for b in bullets[1:]:
            p = body.add_paragraph()
            p.text = b
            for run in p.runs:
                run.font.size = Pt(18)

    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue(), {
        "company": company,
        "audience": audience,
        "slide_count": len(bullets_per_slide) + 1,
        "industry": industry_label,
    }


def _gen_xlsx_budget(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a department budget workbook as XLSX bytes."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"

    department = random.choice(_ctx(industry, "department_names", ["Operations"]))
    fy = random.randint(2024, 2026)

    ws["A1"] = f"{department} — FY{fy} Budget"
    ws["A2"] = f"Prepared by: {fkr.name()}"
    ws["A3"] = f"Last updated: {_now_iso()[:10]}"

    headers = ["Category", "Q1", "Q2", "Q3", "Q4", "Total"]
    for col, h in enumerate(headers, start=1):
        ws.cell(row=5, column=col, value=h)

    # Pool of 16 plausible budget categories — pick a random subset
    # so two budgets for the same FY don't have identical line items.
    full_pool = [
        "Salaries",
        "Benefits",
        "Equipment",
        "Travel",
        "Training",
        "Software",
        "Contractors",
        "Misc",
        "Marketing",
        "Office Lease",
        "Utilities",
        "Cloud Infrastructure",
        "Professional Services",
        "Insurance",
        "Conferences",
        "Recruiting",
    ]
    categories = random.sample(full_pool, k=random.randint(6, len(full_pool)))
    total_rows = 0
    annual_total = 0
    for i, cat in enumerate(categories):
        row = 6 + i
        ws.cell(row=row, column=1, value=cat)
        cat_total = 0
        for q_col in range(2, 6):
            val = random.randint(8000, 250_000)
            ws.cell(row=row, column=q_col, value=val)
            cat_total += val
        ws.cell(row=row, column=6, value=cat_total)
        annual_total += cat_total
        total_rows += 1

    total_row = 6 + len(categories) + 1
    ws.cell(row=total_row, column=1, value="TOTAL")
    for q_col in range(2, 6):
        ws.cell(
            row=total_row,
            column=q_col,
            value=f"=SUM({chr(64 + q_col)}6:{chr(64 + q_col)}{6 + len(categories) - 1})",
        )
    ws.cell(row=total_row, column=6, value=annual_total)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue(), {
        "department": department,
        "fiscal_year": fy,
        "categories": total_rows,
        "annual_total": annual_total,
    }


def _gen_xlsx_inventory(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize an inventory sheet as XLSX bytes."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Inventory"

    # Header set varies — sometimes include warehouse / supplier columns.
    base_headers = [
        "SKU",
        "Description",
        "Category",
        "Qty on Hand",
        "Reorder Point",
        "Unit Cost",
        "Last Counted",
    ]
    optional_headers = ["Warehouse", "Supplier", "Lead Time (days)"]
    optional_chosen = random.sample(optional_headers, k=random.randint(0, len(optional_headers)))
    headers = base_headers + optional_chosen
    for col, h in enumerate(headers, start=1):
        ws.cell(row=1, column=col, value=h)

    categories = _ctx(
        industry, "product_categories", _ctx(industry, "part_categories", ["General"])
    )
    warehouses = [f"WH-{i:02d}" for i in range(1, 9)]
    suppliers = [fkr.company() for _ in range(8)]
    item_count = random.randint(40, 200)
    for i in range(item_count):
        row = 2 + i
        ws.cell(row=row, column=1, value=f"SKU-{random.randint(100000, 999999)}")
        ws.cell(row=row, column=2, value=fkr.bs().capitalize())
        ws.cell(row=row, column=3, value=random.choice(categories))
        ws.cell(row=row, column=4, value=random.randint(0, 5000))
        ws.cell(row=row, column=5, value=random.randint(50, 500))
        ws.cell(row=row, column=6, value=round(random.uniform(2.5, 250.0), 2))
        ws.cell(
            row=row,
            column=7,
            value=fkr.date_between(start_date="-90d", end_date="today").isoformat(),
        )
        # Fill optional columns when present.
        col_offset = 8
        for opt_h in optional_chosen:
            if opt_h == "Warehouse":
                ws.cell(row=row, column=col_offset, value=random.choice(warehouses))
            elif opt_h == "Supplier":
                ws.cell(row=row, column=col_offset, value=random.choice(suppliers))
            elif opt_h == "Lead Time (days)":
                ws.cell(row=row, column=col_offset, value=random.randint(2, 60))
            col_offset += 1

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue(), {
        "item_count": item_count,
        "category_count": len(set(categories)),
        "optional_columns": optional_chosen,
    }


def _gen_eml_message(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize an RFC 5322 email message as .eml bytes.

    Variations: 12 subject templates rotated, AI-drafted body when on,
    optional CC line, optional reply-prefix subject.
    """
    from email.message import EmailMessage

    msg = EmailMessage()
    sender_name = fkr.name()
    sender_email = fkr.email()
    recipient_name = fkr.name()
    recipient_email = fkr.email()
    department = random.choice(_ctx(industry, "department_names", ["Operations"]))
    subjects = [
        f"Q{random.randint(1, 4)} {department} review — agenda",
        "Follow-up on yesterday's call",
        "FYI — updated reporting schedule",
        f"Need your input on the {department.lower()} initiative",
        "Re: budget approval",
        f"Quick question about the {department.lower()} rollout",
        f"Heads-up — {department.lower()} resourcing change",
        f"Action required: {department.lower()} deck review by EOW",
        f"Updated: {department.lower()} risk register",
        f"For review: {department.lower()} forecast revisions",
        f"{department} steering committee — minutes attached",
        f"Re: {department.lower()} contract addendum",
    ]
    subject = random.choice(subjects)
    if _maybe_section(0.3) and not subject.startswith("Re:"):
        subject = "Re: " + subject

    msg["Subject"] = subject
    msg["From"] = f"{sender_name} <{sender_email}>"
    msg["To"] = f"{recipient_name} <{recipient_email}>"
    cc_count = 0
    if _maybe_section(0.4):
        cc_list = [fkr.email() for _ in range(random.randint(1, 3))]
        msg["Cc"] = ", ".join(cc_list)
        cc_count = len(cc_list)
    msg["Date"] = datetime.now(timezone.utc).strftime("%a, %d %b %Y %H:%M:%S +0000")
    msg["Message-ID"] = f"<{uuid.uuid4()}@demo.clone-xs.local>"

    greeting = _rotate(
        f"Hi {recipient_name.split()[0]},",
        f"Hello {recipient_name.split()[0]},",
        f"{recipient_name.split()[0]} —",
    )
    body_text = _maybe_ai(
        ai_client,
        (
            f"Write a 4-6 sentence email body from {sender_name} to {recipient_name} about "
            f"'{subject}' in the {department} group. Tone: collegial business email. End "
            f"with a clear ask. Do NOT include greeting or sign-off."
        ),
        fallback=(
            f"Following up on our discussion about the {department.lower()} workstream. "
            f"I want to make sure we're aligned on the next steps and that nothing falls "
            f"through the cracks before the end of the quarter.\n\n"
            f"Could you confirm a time this week to walk through the latest deck? I "
            f"have a few questions on the rollout timeline and the resource plan."
        ),
        max_tokens=200,
    )
    closing = _rotate("Thanks,", "Best,", "Cheers,", "Talk soon,", "Appreciate it,")

    body = f"{greeting}\n\n{body_text}\n\n{closing}\n{sender_name}\n"
    msg.set_content(body)

    return bytes(msg), {
        "sender": sender_name,
        "recipient": recipient_name,
        "subject": subject,
        "department": department,
        "cc_count": cc_count,
    }


# ── Shared layout helpers for industry-specific generators ──────
#
# The 20 new types share a small set of structural patterns (header
# block + table + footer for PDF; header + body paragraphs + signature
# for DOCX). Hoisting the boilerplate into helpers keeps each generator
# below ~30 lines so the diff stays scannable.


def _pdf_table_doc(
    *,
    title: str,
    header_pairs: list[tuple[str, str]],
    table_header: list[str],
    table_rows: list[list[str]],
    footer_lines: list[str] | None = None,
    pdf_title: str | None = None,
) -> bytes:
    """Build a PDF: title heading, header K/V block, table, footer.

    Used by ~13 of the 20 new generators (account statement, wire
    confirmation, purchase order, receipt, SLA report, BOM, QA report,
    meter reading, transcript, property listing, BOL, customs,
    underwriting report). Each generator only has to compose the
    table data — layout / font / spacing logic lives here.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, title=pdf_title or title)
    styles = getSampleStyleSheet()
    story: list[Any] = [Paragraph(title, styles["Heading2"]), Spacer(1, 8)]
    if header_pairs:
        for k, v in header_pairs:
            story.append(Paragraph(f"<b>{k}:</b> {v}", styles["Normal"]))
        story.append(Spacer(1, 10))
    if table_rows:
        data: list[list[str]] = [table_header, *table_rows]
        tbl = Table(data, hAlign="LEFT")
        tbl.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ]
            )
        )
        story.append(tbl)
        story.append(Spacer(1, 10))
    for line in footer_lines or []:
        story.append(Paragraph(line, styles["Normal"]))
    doc.build(story)
    return buf.getvalue()


def _docx_letter_doc(
    *,
    sender_name: str,
    sender_address: str,
    recipient_name: str,
    recipient_address: str,
    subject: str,
    body_paragraphs: list[str],
    closing: str = "Sincerely,",
) -> bytes:
    """Build a one-page letter-style DOCX. Used for outage notice,
    syllabus, endorsement — anything that's prose-with-headers."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    p = doc.add_paragraph(sender_name)
    p.runs[0].bold = True
    doc.add_paragraph(sender_address)
    doc.add_paragraph()
    doc.add_paragraph(datetime.now().strftime("%B %d, %Y"))
    doc.add_paragraph()
    doc.add_paragraph(recipient_name)
    doc.add_paragraph(recipient_address)
    doc.add_paragraph()
    subj = doc.add_paragraph()
    subj_run = subj.add_run(f"Re: {subject}")
    subj_run.bold = True
    doc.add_paragraph()
    for para in body_paragraphs:
        body_p = doc.add_paragraph(para)
        for run in body_p.runs:
            run.font.size = Pt(11)
    doc.add_paragraph()
    doc.add_paragraph(closing)
    doc.add_paragraph(sender_name)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── New industry-specific generators (20 total) ──────────────────


# Healthcare ──


def _gen_pdf_lab_report(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a clinical lab report as PDF bytes.

    Variations: 3-7 lab panels sampled from a wider pool, AI-drafted
    interpretation footer (60% of the time), rotated department.
    """
    patient = fkr.name()
    mrn = f"MRN-{random.randint(100000, 999999)}"
    ordering_provider = fkr.name()
    accession = f"ACC-{datetime.now().strftime('%Y%m%d')}-{random.randint(1000, 9999)}"
    full_panels = [
        (
            "Hemoglobin",
            f"{random.uniform(11.5, 17.5):.1f} g/dL",
            "12.0–17.0",
            random.choice(["Normal", "Normal", "Normal", "Low"]),
        ),
        ("WBC count", f"{random.uniform(4.0, 11.0):.1f} ×10³/µL", "4.5–11.0", "Normal"),
        (
            "Glucose (fasting)",
            f"{random.randint(75, 130)} mg/dL",
            "70–99",
            random.choice(["Normal", "High"]),
        ),
        (
            "Cholesterol, total",
            f"{random.randint(140, 250)} mg/dL",
            "<200",
            random.choice(["Normal", "Borderline"]),
        ),
        ("Creatinine", f"{random.uniform(0.6, 1.4):.2f} mg/dL", "0.6–1.3", "Normal"),
        (
            "Potassium",
            f"{random.uniform(3.4, 5.2):.2f} mmol/L",
            "3.5–5.1",
            random.choice(["Normal", "Normal", "Low"]),
        ),
        ("Sodium", f"{random.randint(135, 145)} mmol/L", "135–145", "Normal"),
        ("ALT", f"{random.randint(8, 45)} U/L", "7–55", random.choice(["Normal", "High"])),
        ("AST", f"{random.randint(8, 45)} U/L", "8–48", "Normal"),
        (
            "HbA1c",
            f"{round(random.uniform(4.5, 8.5), 1)} %",
            "<5.7",
            random.choice(["Normal", "Pre-diabetic", "Diabetic"]),
        ),
        (
            "Vitamin D, 25-OH",
            f"{random.randint(15, 60)} ng/mL",
            "30–100",
            random.choice(["Normal", "Low"]),
        ),
        ("TSH", f"{round(random.uniform(0.4, 5.0), 2)} mIU/L", "0.4–4.5", "Normal"),
        ("Platelets", f"{random.randint(140, 410)} ×10³/µL", "150–400", "Normal"),
    ]
    panels = random.sample(full_panels, k=random.randint(3, 7))
    abnormal_count = sum(1 for p in panels if p[3] not in ("Normal",))

    department = random.choice(_ctx(industry, "department_names", ["Laboratory"]))
    footer = [
        "<i>Reference ranges are population-based; clinical correlation required.</i>",
        f"Department: {department}",
    ]

    if _maybe_section(0.6):
        notes = _maybe_ai(
            ai_client,
            f"Write a 1-2 sentence clinical interpretation note for a lab report ordered by {ordering_provider} for patient {patient}. {abnormal_count} of {len(panels)} analytes flagged abnormal. Tone: clinical, third-person, brief.",
            fallback=_rotate(
                f"Results reviewed; {abnormal_count} flagged value(s) require clinical correlation with patient history.",
                "All values within expected variability for the patient demographic; no urgent follow-up indicated."
                if abnormal_count == 0
                else "Recommend follow-up testing in 4-6 weeks for the flagged analyte(s).",
                "Report electronically signed; results communicated to ordering provider per protocol.",
            ),
            max_tokens=120,
        )
        footer.append(f"<b>Notes:</b> {notes}")

    pdf = _pdf_table_doc(
        title=_rotate("Clinical Laboratory Report", "Lab Results Report", "Patient Lab Report"),
        header_pairs=[
            ("Patient", patient),
            ("MRN", mrn),
            ("Accession", accession),
            ("Ordering provider", ordering_provider),
            ("Collected", _now_iso()),
        ],
        table_header=["Analyte", "Result", "Reference range", "Flag"],
        table_rows=[list(r) for r in panels],
        footer_lines=footer,
        pdf_title="Lab report",
    )
    return pdf, {
        "patient": patient,
        "mrn": mrn,
        "accession": accession,
        "panels": len(panels),
        "abnormal_count": abnormal_count,
    }


def _gen_pdf_discharge_summary(
    industry: str, fkr: Any, ai_client: Any | None
) -> tuple[bytes, dict]:
    """Synthesize a multi-paragraph patient discharge summary.

    Variations: AI-drafted hospital-course, discharge-meds, and
    follow-up sections; optional 'Activity restrictions' section
    (50% of the time); rotated heading style.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    patient = fkr.name()
    mrn = f"MRN-{random.randint(100000, 999999)}"
    los = random.randint(2, 10)
    admission = (datetime.now() - timedelta(days=los)).date().isoformat()
    discharge = datetime.now().date().isoformat()
    diagnosis = random.choice(_ctx(industry, "claim_diagnoses", ["—"]))
    department = random.choice(_ctx(industry, "department_names", ["Internal Medicine"]))
    title_variant = _rotate(
        "PATIENT DISCHARGE SUMMARY", "DISCHARGE SUMMARY", "INPATIENT DISCHARGE NOTE"
    )

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, title="Discharge summary")
    styles = getSampleStyleSheet()
    story: list[Any] = [
        Paragraph(title_variant, styles["Heading2"]),
        Spacer(1, 8),
        Paragraph(f"<b>Patient:</b> {patient} | <b>MRN:</b> {mrn}", styles["Normal"]),
        Paragraph(
            f"<b>Admission:</b> {admission} | <b>Discharge:</b> {discharge} | <b>LOS:</b> {los}d",
            styles["Normal"],
        ),
        Paragraph(f"<b>Service:</b> {department}", styles["Normal"]),
        Spacer(1, 10),
        Paragraph("<b>Principal diagnosis</b>", styles["Heading4"]),
        Paragraph(diagnosis, styles["Normal"]),
        Spacer(1, 6),
        Paragraph("<b>Hospital course</b>", styles["Heading4"]),
        Paragraph(
            _maybe_ai(
                ai_client,
                f"Write a 3-4 sentence hospital-course narrative for patient with diagnosis '{diagnosis}', length-of-stay {los} days, on the {department} service. Tone: clinical, concise, third-person.",
                fallback=(
                    f"Patient was admitted for evaluation and stabilization of {diagnosis.split('(')[0].strip()}. "
                    f"Treatment proceeded without complication over the {los}-day stay; vital signs remained within "
                    f"expected ranges throughout. The care team coordinated discharge planning with the patient and family."
                ),
                max_tokens=180,
            ),
            styles["Normal"],
        ),
        Spacer(1, 6),
        Paragraph("<b>Discharge medications</b>", styles["Heading4"]),
        Paragraph(
            _maybe_ai(
                ai_client,
                f"Write 2 sentences listing typical discharge-medication instructions following hospitalization for {diagnosis.split('(')[0].strip()}. Tone: clinical.",
                fallback=_rotate(
                    "Continue prior medications. New prescriptions reviewed with the patient; follow-up in 7-10 days with primary care.",
                    "Resume home medications as previously prescribed. Two new prescriptions issued with full counselling on dosing and side-effects.",
                    "Medication reconciliation completed at discharge. Patient verbalised understanding of regimen and follow-up schedule.",
                ),
                max_tokens=120,
            ),
            styles["Normal"],
        ),
        Spacer(1, 6),
        Paragraph("<b>Follow-up</b>", styles["Heading4"]),
        Paragraph(
            _maybe_ai(
                ai_client,
                f"Write 2 sentences of follow-up instructions: scheduling a {department} appointment + return-precaution criteria. Tone: clinical, patient-facing.",
                fallback=(
                    f"Schedule follow-up appointment with {department} within two weeks. "
                    "Return to the emergency department for any new chest pain, shortness of breath, or fever above 38.5°C."
                ),
                max_tokens=140,
            ),
            styles["Normal"],
        ),
    ]
    if _maybe_section(0.5):
        story.append(Spacer(1, 6))
        story.append(Paragraph("<b>Activity restrictions</b>", styles["Heading4"]))
        story.append(
            Paragraph(
                _rotate(
                    "No heavy lifting (>10 lbs) for 2 weeks. May resume normal activity as tolerated.",
                    "Bed rest for the first 48 hours; gradual return to activities thereafter as comfort allows.",
                    "Avoid driving or operating machinery for 72 hours after discharge.",
                    "Light activity only for 7 days; no strenuous exercise until clinical follow-up.",
                ),
                styles["Normal"],
            )
        )
    doc.build(story)
    return buf.getvalue(), {
        "patient": patient,
        "mrn": mrn,
        "diagnosis": diagnosis,
        "length_of_stay_days": los,
        "department": department,
    }


# Financial ──


def _gen_pdf_account_statement(
    industry: str, fkr: Any, ai_client: Any | None
) -> tuple[bytes, dict]:
    """Synthesize a monthly account statement.

    Variations: 5-20 transactions, AI-drafted notices footer (50% of
    the time), rotated statement title.
    """
    holder = fkr.name()
    account_no = f"****{random.randint(1000, 9999)}"
    period = datetime.now().strftime("%B %Y")
    txn_types = _ctx(industry, "transaction_types", ["Transaction"])
    rows = []
    balance = round(random.uniform(2500, 25000), 2)
    txn_count = random.randint(5, 20)
    for i in range(txn_count):
        date = (datetime.now() - timedelta(days=i * random.randint(1, 4))).date().isoformat()
        ttype = random.choice(txn_types)
        amount = round(random.uniform(-500, 1500), 2)
        balance = round(balance + amount, 2)
        rows.append([date, ttype, f"${amount:+,.2f}", f"${balance:,.2f}"])

    footer = [
        "<i>"
        + _rotate(
            "Please review your statement carefully. Report any discrepancy within 60 days.",
            "Statement reflects activity through the period close. Pending transactions excluded.",
            "Notice: federal regulation requires you to report errors within 60 days of the statement date.",
        )
        + "</i>"
    ]
    if _maybe_section(0.5):
        notices = _maybe_ai(
            ai_client,
            f"Write a 2-sentence 'Important notices' paragraph for a {period} bank account statement. Cover one regulatory tip + one fee/rate update. Tone: bank-formal, brief.",
            fallback=_rotate(
                "Important: variable-rate account APY changed effective the 15th of the period. See terms-and-conditions for current schedule.",
                "Notice: enhanced fraud monitoring is now active on all accounts. Set up custom alerts in online banking.",
                "Reminder: paper-statement fee waived for accounts opted into e-delivery. Update preferences in online banking.",
            ),
            max_tokens=140,
        )
        footer.append(f"<b>Notices:</b> {notices}")

    pdf = _pdf_table_doc(
        title=_rotate(
            f"Account Statement — {period}",
            f"Monthly Statement — {period}",
            f"Statement of Account — {period}",
        ),
        header_pairs=[
            ("Account holder", holder),
            ("Account number", account_no),
            ("Statement period", period),
            ("Closing balance", f"${balance:,.2f}"),
        ],
        table_header=["Date", "Description", "Amount", "Balance"],
        table_rows=rows,
        footer_lines=footer,
        pdf_title="Account statement",
    )
    return pdf, {
        "holder": holder,
        "account_no": account_no,
        "transactions": len(rows),
        "closing_balance": balance,
    }


def _gen_pdf_wire_confirmation(
    industry: str, fkr: Any, ai_client: Any | None
) -> tuple[bytes, dict]:
    """Synthesize a single-page wire transfer confirmation.

    Variations: SWIFT pool of 12 banks, optional intermediary-bank
    line (40% of the time), AI-drafted confirmation note.
    """
    sender = fkr.name()
    recipient = fkr.name()
    amount = round(random.uniform(1000, 250000), 2)
    wire_id = f"WIRE-{datetime.now().strftime('%Y%m%d')}-{random.randint(100000, 999999)}"
    swift_pool = [
        "CHASUS33",
        "BARCGB22",
        "DEUTDEFF",
        "CITIUS33",
        "BNPAFRPP",
        "HSBCGB2L",
        "MIDLGB22",
        "WELLUS6S",
        "BOFAUS3N",
        "USBKUS44",
        "PNCCUS33",
        "TDOMCATTTOR",
    ]
    swift = random.choice(swift_pool)
    currency = random.choice(["USD", "USD", "USD", "EUR", "GBP", "CAD"])
    pairs = [
        ("Wire ID", wire_id),
        ("Sender", sender),
        ("Recipient", recipient),
        ("Amount", f"{currency} {amount:,.2f}"),
        ("SWIFT/BIC", swift),
        ("Value date", datetime.now().date().isoformat()),
        ("Status", _rotate("Settled", "Completed", "Confirmed")),
    ]
    if _maybe_section(0.4):
        pairs.insert(5, ("Intermediary bank", random.choice(swift_pool)))

    note = _maybe_ai(
        ai_client,
        f"Write a 1-2 sentence formal confirmation note for a wire transfer of {currency} {amount:.2f} from {sender} to {recipient}, settled today. Tone: bank-formal, brief.",
        fallback=_rotate(
            "This confirmation evidences a completed funds transfer. Retain for your records.",
            "Funds have been irrevocably released to the receiving institution per your instruction. Settlement is final.",
            "Wire executed against good and collected funds. No further action required from your end.",
        ),
        max_tokens=120,
    )

    pdf = _pdf_table_doc(
        title="Wire Transfer Confirmation",
        header_pairs=pairs,
        table_header=[],
        table_rows=[],
        footer_lines=[
            f"<i>{note}</i>",
            f"Reference: {fkr.bothify(text='REF-?##??##').upper()}",
        ],
        pdf_title="Wire confirmation",
    )
    return pdf, {
        "wire_id": wire_id,
        "amount": amount,
        "currency": currency,
        "swift": swift,
    }


# Retail / manufacturing ──


def _gen_pdf_purchase_order(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a purchase order with line items.

    Variations: 3-12 line items, AI-drafted special-instructions
    paragraph (50% of the time), variable lead time, optional
    incoterm.
    """
    buyer = fkr.company()
    supplier = fkr.company()
    po_number = f"PO-{datetime.now().strftime('%Y%m')}-{random.randint(10000, 99999)}"
    cats = _ctx(industry, "product_categories", _ctx(industry, "part_categories", ["Item"]))
    incoterms = _ctx("logistics", "incoterms", ["FOB"])
    rows = []
    total = 0.0
    line_count = random.randint(3, 12)
    for _ in range(line_count):
        item = random.choice(cats)
        qty = random.randint(10, 500)
        unit = round(random.uniform(5, 250), 2)
        line_total = round(qty * unit, 2)
        total += line_total
        rows.append([item, str(qty), f"${unit:,.2f}", f"${line_total:,.2f}"])
    rows.append(["", "", "Total", f"${total:,.2f}"])

    lead_days = random.choice([7, 14, 21, 30, 45, 60])
    pairs = [
        ("Buyer", buyer),
        ("Supplier", supplier),
        ("PO number", po_number),
        ("Order date", datetime.now().date().isoformat()),
        ("Required by", (datetime.now() + timedelta(days=lead_days)).date().isoformat()),
    ]
    if _maybe_section(0.5):
        pairs.append(("Incoterm", random.choice(incoterms)))

    footer = [
        "<i>"
        + _rotate(
            "All goods subject to acceptance inspection on receipt.",
            "Goods must comply with attached specifications. Non-conforming items returned at supplier expense.",
            "Acceptance is contingent on compliance with the master purchase agreement and applicable specifications.",
        )
        + "</i>"
    ]
    if _maybe_section(0.5):
        instr = _maybe_ai(
            ai_client,
            f"Write a 1-2 sentence 'Special instructions' note for purchase order {po_number} from {buyer} to {supplier}, total ${total:,.2f}, lead time {lead_days} days. Tone: procurement-formal, brief.",
            fallback=_rotate(
                "Confirm shipment within 48 hours; expedite if any line cannot ship by the required date.",
                "Pack each line on its own pallet. Include packing list and certificate of conformance with each shipment.",
                "Partial shipments not permitted; complete order to ship in a single dispatch.",
            ),
            max_tokens=120,
        )
        footer.append(f"<b>Special instructions:</b> {instr}")

    pdf = _pdf_table_doc(
        title=_rotate(
            f"Purchase Order — {po_number}",
            f"PO {po_number} — Goods Order",
            f"Order Confirmation — {po_number}",
        ),
        header_pairs=pairs,
        table_header=["Item", "Qty", "Unit price", "Line total"],
        table_rows=rows,
        footer_lines=footer,
        pdf_title="Purchase order",
    )
    return pdf, {
        "po_number": po_number,
        "buyer": buyer,
        "supplier": supplier,
        "line_count": line_count,
        "lead_days": lead_days,
        "total_usd": round(total, 2),
    }


def _gen_pdf_receipt(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a point-of-sale sales receipt.

    Variations: 1-9 line items, varying tax rate, AI-drafted thank-you
    message (40% of the time), optional loyalty-points footer.
    """
    store = random.choice(_ctx(industry, "store_codes", ["STR-1001"]))
    txn_id = f"TXN-{datetime.now().strftime('%Y%m%d%H%M%S')}-{random.randint(100, 999)}"
    items = []
    subtotal = 0.0
    line_count = random.randint(1, 9)
    for _ in range(line_count):
        name = fkr.word().title()
        qty = random.randint(1, 4)
        price = round(random.uniform(2.99, 79.99), 2)
        items.append([name, str(qty), f"${price * qty:.2f}"])
        subtotal += price * qty
    tax_rate = random.choice([0.0625, 0.0775, 0.08875, 0.10, 0.0925])
    tax = round(subtotal * tax_rate, 2)
    total = round(subtotal + tax, 2)
    items.append(["", "Subtotal", f"${subtotal:.2f}"])
    items.append(["", f"Sales tax ({tax_rate * 100:.3f}%)", f"${tax:.2f}"])
    items.append(["", "TOTAL", f"${total:.2f}"])

    footer = [
        "<i>"
        + _maybe_ai(
            ai_client,
            f"Write a 1-sentence cheerful customer thank-you message for a {industry} retail receipt, total ${total:.2f}. Mention return policy.",
            fallback=_rotate(
                "Thank you for your purchase. Returns accepted within 30 days with original receipt.",
                "Thanks for shopping with us! Hold onto your receipt — returns within 60 days, no questions asked.",
                "We appreciate your business. Returns and exchanges welcomed within 30 days.",
                "Thank you for being a valued customer. Visit us online for exclusive offers.",
            ),
            max_tokens=80,
        )
        + "</i>"
    ]
    if _maybe_section(0.5):
        footer.append(f"Loyalty points earned: {int(total)}")

    pdf = _pdf_table_doc(
        title=_rotate("Sales Receipt", "Customer Receipt", "Purchase Receipt"),
        header_pairs=[
            ("Store", store),
            ("Transaction", txn_id),
            ("Cashier", fkr.first_name()),
            ("Date", _now_iso()),
        ],
        table_header=["Item", "Qty", "Price"],
        table_rows=items,
        footer_lines=footer,
        pdf_title="Sales receipt",
    )
    return pdf, {
        "store": store,
        "transaction": txn_id,
        "total_usd": total,
        "tax_rate": tax_rate,
        "line_count": line_count,
    }


# Telecom / energy ──


def _gen_pdf_sla_report(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a monthly SLA-attainment report.

    Variations: 3-8 metrics sampled from the (now expanded) pool,
    AI-drafted exec summary + recommendations footer, varying
    breach distribution.
    """
    customer = fkr.company()
    period = datetime.now().strftime("%B %Y")
    full_pool = _ctx(industry, "sla_metrics", ["Network availability ≥ 99.9%"])
    metrics = random.sample(full_pool, k=min(random.randint(3, 8), len(full_pool)))
    rows = []
    breaches = 0
    for m in metrics:
        attained = random.uniform(99.0, 99.99)
        # Parse the trailing "≥ NN%" — tolerant of any extra words after
        # the percent (e.g. "≥ 92% of footprint"). Falls back to 99.5%
        # when the metric isn't a percentage SLA.
        target = 99.5
        if "≥" in m:
            tail = m.rsplit("≥", 1)[-1].strip()
            digits = ""
            for ch in tail:
                if ch.isdigit() or ch == ".":
                    digits += ch
                elif digits:
                    break
            try:
                target = float(digits) if digits else 99.5
            except ValueError:
                target = 99.5
        status = "Met" if attained >= target else "Breach"
        if status == "Breach":
            breaches += 1
        rows.append([m, f"{attained:.2f}%", status])

    summary = _maybe_ai(
        ai_client,
        f"Write a 2-sentence executive summary for the {period} SLA report for customer {customer}: {len(metrics)} metrics tracked, {breaches} breaches. Tone: account-management, brief.",
        fallback=_rotate(
            f"Performance for {period} {'met all SLA commitments' if breaches == 0 else f'showed {breaches} breach(es)'} across {len(metrics)} metrics. Service credits {'do not apply' if breaches == 0 else 'will be calculated per the MSA'}.",
            f"{period} closes with {breaches} SLA breach(es) out of {len(metrics)} tracked metrics. We are reviewing root cause and corrective actions with engineering leadership.",
            f"Service performance in {period} was within commitment thresholds on all but {breaches} metric(s); detailed root-cause review is in progress.",
        ),
        max_tokens=140,
    )

    pdf = _pdf_table_doc(
        title=_rotate(
            f"Service Level Agreement Report — {period}",
            f"SLA Performance Report — {period}",
            f"Service Performance Report — {period}",
        ),
        header_pairs=[
            ("Customer", customer),
            ("Reporting period", period),
            ("Total breaches", str(breaches)),
            ("Account manager", fkr.name()),
        ],
        table_header=["Metric", "Attained", "Status"],
        table_rows=rows,
        footer_lines=[
            f"<b>Executive summary:</b> {summary}",
            "<i>Service credits, where applicable, are calculated per the master service agreement and applied to the next invoice.</i>",
        ],
        pdf_title="SLA report",
    )
    return pdf, {
        "customer": customer,
        "period": period,
        "metrics_tracked": len(rows),
        "breaches": breaches,
    }


def _gen_docx_outage_notice(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a service-outage notification letter.

    Variations: AI-drafted explanation paragraph, optional credit-
    offer paragraph (40% of the time), variable duration, rotated
    closing.
    """
    cause = random.choice(_ctx(industry, "outage_causes", ["Equipment failure"]))
    affected = random.randint(50, 50000)
    duration_min = random.choice([15, 30, 45, 60, 90, 120, 180, 240])
    incident_id = f"INC-{datetime.now().strftime('%Y%m%d')}-{random.randint(100, 999)}"
    customer = fkr.name()
    incident_date = datetime.now().strftime("%B %d, %Y")

    paragraphs = [
        _maybe_ai(
            ai_client,
            f"Write a 2-sentence opening for a customer outage notification. Date: {incident_date}, root cause: {cause}, duration: {duration_min} minutes. Tone: formal, apologetic, factual.",
            fallback=(
                f"We are writing to inform you of a service-affecting incident that occurred on {incident_date}. "
                f"The root cause has been identified as: {cause}. Total duration: {duration_min} minutes."
            ),
            max_tokens=130,
        ),
        _maybe_ai(
            ai_client,
            f"Write a 2-sentence paragraph describing impact and remediation: ~{affected:,} customers affected, restored after engineering remediation. Tone: neutral.",
            fallback=(
                f"Approximately {affected:,} customers were impacted in the affected service area. "
                "Service was fully restored after our field engineering team completed the necessary remediation."
            ),
            max_tokens=130,
        ),
        _rotate(
            "We take service reliability seriously, and we apologize for the inconvenience this incident may have caused. A detailed root-cause analysis will follow within 5 business days.",
            "We sincerely apologize for the disruption. Our post-incident review is underway and the full root-cause analysis will be shared with you within one week.",
            "Please accept our apologies for the impact this incident had on your operations. We are committed to preventing recurrence and will share the post-mortem within five business days.",
        ),
        f"If you experience any continuing issues, please reference incident {incident_id} when contacting our support team.",
    ]
    if _maybe_section(0.4):
        credit_pct = random.choice([5, 10, 15, 25])
        paragraphs.insert(
            3,
            _maybe_ai(
                ai_client,
                f"Write a 1-sentence service-credit offer paragraph: {credit_pct}% credit on the next invoice as a goodwill gesture. Tone: formal.",
                fallback=f"As a goodwill gesture, we will apply a {credit_pct}% service credit to your next invoice; no action is required on your part.",
                max_tokens=80,
            ),
        )

    docx = _docx_letter_doc(
        sender_name=_rotate(
            "Network Operations Center", "Service Reliability Team", "Customer Operations"
        ),
        sender_address=f"{fkr.company()}\n{fkr.address()}",
        recipient_name=customer,
        recipient_address=fkr.address(),
        subject=f"Service incident notification — {incident_id}",
        body_paragraphs=paragraphs,
        closing=_rotate("Regards,", "Sincerely,", "With apologies,", "Respectfully,"),
    )
    return docx, {
        "incident_id": incident_id,
        "cause": cause,
        "affected_customers": affected,
        "duration_min": duration_min,
        "paragraph_count": len(paragraphs),
    }


# Manufacturing ──


def _gen_pdf_bom(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a multi-line bill of materials.

    Variations: 6-25 line items, AI-drafted engineering-notes footer
    (50% of the time), revision letter from a wider pool.
    """
    assembly = f"ASM-{random.randint(10000, 99999)}"
    revision = f"Rev {random.choice(['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H'])}"
    parts = _ctx(industry, "part_categories", ["Component"])
    rows = []
    line_count = random.randint(6, 25)
    total_cost = 0.0
    for i in range(line_count):
        part_id = f"P-{random.randint(10000, 99999)}"
        cat = random.choice(parts)
        qty = random.randint(1, 12)
        unit_cost = round(random.uniform(0.50, 350.00), 2)
        total_cost += qty * unit_cost
        rows.append([str(i + 1), part_id, cat, str(qty), f"${unit_cost:,.2f}"])

    footer = [
        "<i>"
        + _rotate(
            "Substitutions require engineering change order approval.",
            "All parts subject to incoming-inspection criteria per IQS-001.",
            "Sourcing changes must be approved by Manufacturing Engineering before production.",
        )
        + "</i>"
    ]
    if _maybe_section(0.5):
        notes = _maybe_ai(
            ai_client,
            f"Write a 1-2 sentence engineering note for BOM {assembly} {revision} ({line_count} line items, ${total_cost:.2f} total). Tone: technical, brief.",
            fallback=_rotate(
                "Lead-time critical: hydraulic components have a 12-week typical lead time and must be ordered ahead of kit-out.",
                "Configuration deviates from previous revision; review with Manufacturing Engineering before kicking off the run.",
                "All fasteners must conform to Grade 8.8 specification; reference SOP-PR-014 for inspection procedure.",
            ),
            max_tokens=120,
        )
        footer.append(f"<b>Engineering notes:</b> {notes}")

    pdf = _pdf_table_doc(
        title=f"Bill of Materials — {assembly} ({revision})",
        header_pairs=[
            ("Assembly", assembly),
            ("Revision", revision),
            ("Released by", fkr.name()),
            ("Released", datetime.now().date().isoformat()),
            ("Total parts", str(line_count)),
            ("Estimated cost", f"${total_cost:,.2f}"),
        ],
        table_header=["#", "Part ID", "Category", "Qty", "Unit cost"],
        table_rows=rows,
        footer_lines=footer,
        pdf_title="Bill of materials",
    )
    return pdf, {
        "assembly": assembly,
        "revision": revision,
        "line_count": line_count,
        "total_cost": total_cost,
    }


def _gen_pdf_qa_report(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a quality inspection report.

    Variations: 4-9 inspection checks sampled from a wider pool,
    AI-drafted findings paragraph (60% of the time), variable
    sample size.
    """
    lot = f"LOT-{datetime.now().strftime('%Y%m%d')}-{random.randint(1000, 9999)}"
    inspector = fkr.name()
    sampled = random.randint(50, 1000)
    failures = random.randint(0, max(1, sampled // 20))
    full_checks = [
        ("Dimensional accuracy (±0.05mm)", random.choice(["Pass", "Pass", "Pass", "Fail"])),
        ("Surface finish (Ra ≤ 1.6µm)", random.choice(["Pass", "Pass", "Pass"])),
        ("Material certification", "Pass"),
        ("Visual inspection — defects", random.choice(["Pass", "Pass", "Fail"])),
        ("Functional test", random.choice(["Pass", "Pass"])),
        ("Hardness test (HRC)", random.choice(["Pass", "Pass", "Pass"])),
        ("Torque verification", random.choice(["Pass", "Pass"])),
        ("Coating thickness", random.choice(["Pass", "Pass", "Fail"])),
        ("Pull test (mounting)", random.choice(["Pass", "Pass", "Pass"])),
        ("Weight tolerance", "Pass"),
        ("Electrical continuity", random.choice(["Pass", "Pass"])),
    ]
    checks = random.sample(full_checks, k=random.randint(4, 9))
    rows = [[c, r, fkr.bothify(text="DR-####")] for c, r in checks]
    fail_count = sum(1 for _, r in checks if r == "Fail")

    footer = [
        "<i>"
        + _rotate(
            "Lots with any failed checks are held pending engineering review.",
            "Disposition is final pending MRB review per QMS-002.",
            "All deviations are tracked in the Quality Management System; full traceability available on request.",
        )
        + "</i>"
    ]
    if _maybe_section(0.6):
        findings = _maybe_ai(
            ai_client,
            f"Write 2 sentences of QA findings narrative for lot {lot}: {sampled} units sampled, {failures} failures, {fail_count} of {len(checks)} check categories failed. Tone: quality-assurance technical.",
            fallback=_rotate(
                f"Sample of {sampled} units yielded {failures} failures; defect mode primarily relates to dimensional drift on the OD turning operation.",
                f"{fail_count} of {len(checks)} check categories returned Fail dispositions; corrective actions opened with the responsible production cell.",
                f"Inspection complete. Of {sampled} units, {failures} did not meet acceptance criteria. Recommend tightening process controls on the affected operation.",
            ),
            max_tokens=140,
        )
        footer.append(f"<b>Findings:</b> {findings}")

    pdf = _pdf_table_doc(
        title=f"Quality Inspection Report — {lot}",
        header_pairs=[
            ("Lot number", lot),
            ("Inspector", inspector),
            ("Inspection date", datetime.now().date().isoformat()),
            ("Sample size", str(sampled)),
            ("Failures", str(failures)),
            ("Disposition", "Released" if failures == 0 else "Hold for review"),
        ],
        table_header=["Check", "Result", "Deviation report"],
        table_rows=rows,
        footer_lines=footer,
        pdf_title="QA report",
    )
    return pdf, {
        "lot": lot,
        "sampled": sampled,
        "failures": failures,
        "check_count": len(checks),
        "fail_categories": fail_count,
    }


def _gen_pdf_meter_reading(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a multi-meter reading report (utility / energy).

    Variations: 5-25 meters, varying rate per kWh, AI-drafted period
    summary footer (40% of the time).
    """
    period = datetime.now().strftime("%B %Y")
    rate = round(random.uniform(0.10, 0.22), 4)
    rows = []
    total_delta = 0
    meter_count = random.randint(5, 25)
    for _ in range(meter_count):
        meter_id = f"M-{random.randint(100000, 999999)}"
        reading = random.randint(8000, 25000)
        delta = random.randint(200, 2400)
        total_delta += delta
        rows.append([meter_id, str(reading), f"+{delta}", f"${delta * rate:.2f}"])

    footer = [
        "<i>"
        + _rotate(
            "All readings verified to ±2 kWh per applicable utility commission rules.",
            "Readings collected via AMI; manual verification on a sampled basis.",
            "Estimated readings flagged for next-cycle physical verification.",
        )
        + "</i>"
    ]
    if _maybe_section(0.4):
        summary = _maybe_ai(
            ai_client,
            f"Write a 1-sentence period summary for {period} meter readings: {meter_count} meters, total {total_delta} kWh consumed, rate ${rate}/kWh. Tone: utility-formal, brief.",
            fallback=_rotate(
                f"Aggregate consumption for the {period} cycle was {total_delta:,} kWh across {meter_count} services, in line with seasonal expectations.",
                f"Total billable kWh for {period}: {total_delta:,}, charged at ${rate:.4f}/kWh per the current tariff schedule.",
                f"{period} period closed with {meter_count} meters read and {total_delta:,} kWh consumed; no anomalies flagged.",
            ),
            max_tokens=100,
        )
        footer.append(f"<b>Period summary:</b> {summary}")

    pdf = _pdf_table_doc(
        title=f"Meter Reading Report — {period}",
        header_pairs=[
            ("Reading period", period),
            ("Service area", fkr.city()),
            ("Read by", fkr.name()),
            ("Tariff rate", f"${rate:.4f}/kWh"),
        ],
        table_header=["Meter ID", "Reading (kWh)", "Δ from prior", "Charge"],
        table_rows=rows,
        footer_lines=footer,
        pdf_title="Meter reading",
    )
    return pdf, {
        "period": period,
        "meters_read": meter_count,
        "total_kwh": total_delta,
        "rate_per_kwh": rate,
    }


# Education ──


def _gen_pdf_transcript(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize an academic transcript.

    Variations: 4-12 courses sampled from the wider pool, AI-drafted
    honors / academic standing footer (40% of the time), program
    pool of 12 degrees.
    """
    student = fkr.name()
    student_id = f"S-{random.randint(1000000, 9999999)}"
    program = random.choice(
        [
            "B.Sc. Computer Science",
            "B.A. English Literature",
            "B.S. Biology",
            "M.Sc. Data Science",
            "B.A. History",
            "B.S. Mechanical Engineering",
            "B.A. Economics",
            "M.B.A.",
            "B.S. Mathematics",
            "B.A. Psychology",
            "B.S. Chemistry",
            "M.A. Public Policy",
        ]
    )
    full_courses = _ctx(industry, "course_codes", ["CS101 — Intro"])
    courses = random.sample(full_courses, k=min(random.randint(4, 12), len(full_courses)))
    grades = _ctx(industry, "grade_letters", ["A", "B", "C"])
    rows = []
    gpa_total = 0.0
    grade_pts = {
        "A": 4.0,
        "A-": 3.7,
        "B+": 3.3,
        "B": 3.0,
        "B-": 2.7,
        "C+": 2.3,
        "C": 2.0,
        "C-": 1.7,
        "D+": 1.3,
        "D": 1.0,
        "P": None,
        "NP": None,
        "I": None,
        "W": None,
        "AU": None,
    }
    for c in courses:
        g = random.choice(grades)
        credits = random.choice([2, 3, 3, 3, 4, 4, 5])
        rows.append([c, str(credits), g])
        if grade_pts.get(g) is not None:
            gpa_total += grade_pts[g] * credits  # type: ignore[operator]
    total_credits = sum(int(r[1]) for r in rows)
    gpa = gpa_total / total_credits if total_credits else 0.0

    standing = (
        "Dean's List" if gpa >= 3.7 else ("Good standing" if gpa >= 2.5 else "Academic probation")
    )
    footer = [
        "<i>"
        + _rotate(
            "This transcript is issued by the Office of the Registrar. Tampering invalidates the document.",
            "Official transcript bearing the seal of the Registrar. Unauthorized alteration is prohibited.",
            "This document is the official record of the student's academic performance. Replicas without seal are unofficial.",
        )
        + "</i>"
    ]
    if _maybe_section(0.4):
        narrative = _maybe_ai(
            ai_client,
            f"Write a 1-sentence honors / academic-standing note for {student} in program '{program}', GPA {gpa:.2f}, currently {standing}. Tone: registrar-formal.",
            fallback=_rotate(
                f"Student is in {standing} as of the most recent term.",
                f"Cumulative academic standing: {standing}. Eligible for full enrollment in subsequent terms.",
                f"Student has met the program's academic progress requirements; standing: {standing}.",
            ),
            max_tokens=80,
        )
        footer.append(f"<b>Standing:</b> {narrative}")

    pdf = _pdf_table_doc(
        title=_rotate(
            "Official Academic Transcript", "Student Academic Record", "Transcript of Record"
        ),
        header_pairs=[
            ("Student", student),
            ("Student ID", student_id),
            ("Program", program),
            ("Cumulative GPA", f"{gpa:.2f}"),
            ("Total credits", str(total_credits)),
        ],
        table_header=["Course", "Credits", "Grade"],
        table_rows=rows,
        footer_lines=footer,
        pdf_title="Transcript",
    )
    return pdf, {
        "student": student,
        "student_id": student_id,
        "program": program,
        "gpa": round(gpa, 2),
        "standing": standing,
        "courses": len(rows),
        "total_credits": total_credits,
    }


def _gen_docx_syllabus(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a course syllabus DOCX.

    Variations: AI-drafted course-objectives + grading-policy
    paragraphs, randomized assessment weights, optional textbook
    paragraph (50% of the time).
    """
    course = random.choice(_ctx(industry, "course_codes", ["CS101 — Intro"]))
    instructor = fkr.name()
    term = random.choice(_ctx(industry, "academic_terms", ["Fall 2025"]))
    department = random.choice(_ctx(industry, "department_names", ["Academic Affairs"]))

    # Random assessment weights that sum to 100.
    ps_w = random.choice([20, 25, 30, 35])
    mid_w = random.choice([20, 25, 30])
    final_w = random.choice([25, 30, 35])
    part_w = 100 - ps_w - mid_w - final_w

    paragraphs = [
        f"Welcome to {course}, taught by {instructor} for the {term} term. This syllabus outlines the course goals, weekly schedule, and grading policy.",
        _maybe_ai(
            ai_client,
            f"Write a 2-sentence course-objectives paragraph for {course} taught by {instructor}. Industry context: {industry}. Tone: academic, brief.",
            fallback="Course objectives: students will develop a working knowledge of the core concepts, complete weekly problem sets, and participate in a term project that synthesizes the material.",
            max_tokens=140,
        ),
        f"Assessment: {ps_w}% problem sets, {mid_w}% midterm examination, {final_w}% final project, {max(part_w, 0)}% class participation. Late work is penalized 10% per day unless prior arrangement is made.",
        _rotate(
            "Office hours: Tuesdays and Thursdays, 2:00–3:30 PM, in the instructor's office.",
            "Office hours: Mondays 10:00–11:30 AM and Fridays 1:00–2:30 PM, in person or via Zoom.",
            "Office hours by appointment via the LMS scheduler. In-person sessions held in the department office.",
        )
        + " Email is the preferred contact method for scheduling outside office hours.",
        _maybe_ai(
            ai_client,
            f"Write a 1-2 sentence academic-integrity policy paragraph for syllabus of {course}. Tone: academic, firm but not threatening.",
            fallback="Academic integrity: all work must be your own except where explicitly designated as group work. Suspected violations are referred to the Office of Academic Integrity.",
            max_tokens=120,
        ),
    ]
    if _maybe_section(0.5):
        paragraphs.append(
            _rotate(
                "Required text: as published on the LMS course page; supplementary readings circulated weekly.",
                "Course materials are open-access; readings posted to the LMS at least one week prior to use.",
                "Textbook: details in the welcome packet. PDF copies of each weekly chapter will be made available via the LMS.",
            )
        )

    docx = _docx_letter_doc(
        sender_name=f"Department of {department}",
        sender_address=fkr.company(),
        recipient_name="Enrolled Students",
        recipient_address=term,
        subject=f"Syllabus — {course}",
        body_paragraphs=paragraphs,
        closing=_rotate(
            "Best,", "Looking forward to a great term,", "Welcome aboard,", "Best regards,"
        ),
    )
    return docx, {
        "course": course,
        "instructor": instructor,
        "term": term,
        "department": department,
        "assessment_weights": {
            "problem_sets": ps_w,
            "midterm": mid_w,
            "final": final_w,
            "participation": max(part_w, 0),
        },
        "paragraph_count": len(paragraphs),
    }


# Real estate ──


def _gen_pdf_property_listing(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a property listing one-pager.

    Variations: AI-drafted property description as headline footer
    (60% of the time), variable feature set (3-9 features), wider
    pool of heating/cooling/parking variants.
    """
    address = fkr.address().replace("\n", ", ")
    price = random.randint(250_000, 5_000_000)
    bedrooms = random.choice([1, 2, 3, 3, 4, 4, 5, 6])
    bathrooms = random.choice([1, 1.5, 2, 2.5, 3, 3.5, 4])
    sqft = random.randint(800, 6500)
    year_built = random.randint(1920, 2024)
    listing_id = f"MLS-{random.randint(1000000, 9999999)}"
    ptype = random.choice(_ctx(industry, "property_types", ["Single-family residence"]))
    agent = fkr.name()

    full_features = [
        ["Year built", str(year_built)],
        ["Lot size", f"{random.randint(2500, 25000):,} sqft"],
        [
            "Heating",
            random.choice(
                [
                    "Forced air",
                    "Radiant",
                    "Heat pump",
                    "Baseboard electric",
                    "Geothermal",
                    "Wood stove + central",
                ]
            ),
        ],
        [
            "Cooling",
            random.choice(
                [
                    "Central A/C",
                    "Mini-split",
                    "Window units",
                    "None",
                    "Whole-house fan",
                    "High-efficiency heat pump",
                ]
            ),
        ],
        [
            "Parking",
            random.choice(
                [
                    "2-car attached garage",
                    "1-car detached garage",
                    "Carport",
                    "Street",
                    "Tandem 2-car",
                    "Underground",
                    "Driveway only",
                ]
            ),
        ],
        ["HOA dues", f"${random.choice([0, 0, 0, 150, 250, 425, 600, 850])}/month"],
        [
            "Roof",
            random.choice(
                [
                    "Asphalt shingle (5 yrs)",
                    "Metal (12 yrs)",
                    "Tile (recently inspected)",
                    "Slate (original)",
                ]
            ),
        ],
        [
            "Flooring",
            random.choice(
                ["Hardwood throughout", "Mixed hardwood / carpet", "Tile + carpet", "LVT + carpet"]
            ),
        ],
        [
            "Appliances",
            random.choice(
                [
                    "Stainless steel, included",
                    "Gas range, included",
                    "Dishwasher + fridge included",
                    "All-new (2024)",
                ]
            ),
        ],
        ["School district", fkr.last_name() + " Unified"],
        ["Walkability", f"{random.randint(20, 95)}/100"],
    ]
    features = random.sample(full_features, k=random.randint(3, 9))

    footer = [
        "<i>Information deemed reliable but not guaranteed. Buyer should verify all measurements and details independently.</i>"
    ]
    if _maybe_section(0.6):
        description = _maybe_ai(
            ai_client,
            f"Write a 2-3 sentence enthusiastic property listing description for a {bedrooms}-bedroom {ptype.lower()} at {address}, {sqft} sqft, built {year_built}, list price ${price:,}. Tone: real-estate marketing, evocative.",
            fallback=_rotate(
                f"Beautifully maintained {ptype.lower()} offering {sqft:,} sqft of comfortable living space across {bedrooms} bedrooms. Move-in ready with thoughtful updates throughout. Won't last long at this price.",
                f"Charming {ptype.lower()} in a sought-after neighborhood. Spacious {bedrooms}-bedroom layout with {bathrooms} bathrooms, ideal for entertaining and everyday living. Schedule a private tour today.",
                f"Stunning {bedrooms}-bedroom property with {sqft:,} sqft of designer finishes. Recent updates throughout. A rare find in today's market — submit your best offer.",
            ),
            max_tokens=160,
        )
        footer.insert(0, f"<b>Description:</b> {description}")

    pdf = _pdf_table_doc(
        title=_rotate(
            f"Property Listing — {ptype}", f"For Sale: {ptype}", f"Featured Listing — {ptype}"
        ),
        header_pairs=[
            ("Address", address),
            ("Listing ID", listing_id),
            ("Price", f"${price:,}"),
            ("Bedrooms", str(bedrooms)),
            ("Bathrooms", str(bathrooms)),
            ("Square feet", f"{sqft:,}"),
            ("Listing agent", agent),
        ],
        table_header=["Feature", "Detail"],
        table_rows=features,
        footer_lines=footer,
        pdf_title="Property listing",
    )
    return pdf, {
        "listing_id": listing_id,
        "price_usd": price,
        "bedrooms": bedrooms,
        "bathrooms": bathrooms,
        "sqft": sqft,
        "year_built": year_built,
        "feature_count": len(features),
    }


def _gen_pdf_disclosure(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a property disclosure form.

    Variations: 5-12 disclosure items sampled from the wider pool,
    AI-drafted notes for affirmative answers (when AI mode on),
    rotated title.
    """
    address = fkr.address().replace("\n", ", ")
    seller = fkr.name()
    full_items = _ctx(industry, "disclosure_items", ["—"])
    items = random.sample(full_items, k=min(random.randint(5, 12), len(full_items)))
    rows = []
    affirmative_count = 0
    for item in items:
        ans = random.choice(["Yes", "No", "No", "No", "Unknown"])
        if ans == "Yes":
            affirmative_count += 1
            note = _maybe_ai(
                ai_client,
                f"Write a 1-sentence brief disclosure note explaining a 'Yes' answer to: '{item}'. Tone: factual, real-estate-disclosure.",
                fallback=_rotate(
                    "See attached.",
                    "Documentation provided in addendum.",
                    "Repaired by licensed contractor; receipts available.",
                    "Disclosed per state requirements.",
                ),
                max_tokens=70,
            )
        else:
            note = random.choice(["", "", "N/A"])
        rows.append([item, ans, note])

    pdf = _pdf_table_doc(
        title=_rotate(
            "Seller Property Disclosure",
            "Seller Disclosure Statement",
            "Property Condition Disclosure",
        ),
        header_pairs=[
            ("Property", address),
            ("Seller", seller),
            ("Date", datetime.now().date().isoformat()),
            ("Items disclosed", str(len(rows))),
            ("Affirmative answers", str(affirmative_count)),
        ],
        table_header=["Disclosure item", "Aware of?", "Notes"],
        table_rows=rows,
        footer_lines=[
            "<i>"
            + _rotate(
                "Seller affirms the answers above are true to the best of their knowledge as of the date signed.",
                "All answers are given in good faith based on Seller's actual knowledge as of the signing date.",
                "Seller has not knowingly omitted any material fact regarding the property's condition.",
            )
            + "</i>",
            f"Signed: {seller}    Date: {datetime.now().date().isoformat()}",
        ],
        pdf_title="Property disclosure",
    )
    return pdf, {
        "address": address,
        "seller": seller,
        "items_disclosed": len(rows),
        "affirmative_answers": affirmative_count,
    }


# Logistics ──


def _gen_pdf_bol(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a bill of lading.

    Variations: 2-7 line items, AI-drafted special-handling note
    (45% of the time), wider carrier pool, optional hazmat flag.
    """
    shipper = fkr.company()
    consignee = fkr.company()
    bol_no = f"BOL-{datetime.now().strftime('%Y%m%d')}-{random.randint(10000, 99999)}"
    classes = _ctx(industry, "freight_classes", ["Class 100"])
    carrier_pool = [
        "XPO",
        "Old Dominion",
        "Saia",
        "FedEx Freight",
        "Estes",
        "ABF Freight",
        "TForce Freight",
        "R+L Carriers",
        "YRC",
        "Roadrunner",
    ]
    rows = []
    line_count = random.randint(2, 7)
    total_weight = 0
    for _ in range(line_count):
        descr = fkr.bs().title()
        pkgs = random.randint(1, 30)
        weight = random.randint(50, 8000)
        total_weight += weight
        cls = random.choice(classes)
        rows.append([descr, str(pkgs), f"{weight:,} lbs", cls])

    pairs = [
        ("Shipper", shipper),
        ("Consignee", consignee),
        ("Origin", fkr.city()),
        ("Destination", fkr.city()),
        ("Carrier", random.choice(carrier_pool)),
        ("BOL number", bol_no),
        ("Total weight", f"{total_weight:,} lbs"),
    ]
    is_hazmat = _maybe_section(0.2)
    if is_hazmat:
        pairs.append(
            (
                "Hazmat",
                random.choice(
                    [
                        "Class 3 — Flammable liquid",
                        "Class 8 — Corrosive",
                        "Class 9 — Misc dangerous goods",
                    ]
                ),
            )
        )

    footer = [
        "<i>"
        + _rotate(
            "Received in apparent good order, except as noted, subject to the classifications and rules in effect on the date of issue.",
            "Subject to the classifications and tariffs in effect on the date of issue. Goods inspected and accepted at point of origin.",
            "Goods received in apparent external good order. Carrier liable per Carmack Amendment terms.",
        )
        + "</i>"
    ]
    if _maybe_section(0.45):
        instructions = _maybe_ai(
            ai_client,
            f"Write a 1-2 sentence 'Special handling' note for BOL {bol_no} — {line_count} line items, total {total_weight} lbs{', hazmat shipment' if is_hazmat else ''}. Tone: logistics-formal, brief.",
            fallback=_rotate(
                "Handle with care. Do not stack pallets above one row.",
                "Temperature-sensitive cargo. Maintain ambient between 35-75°F throughout transit.",
                "Time-definite delivery required by end-of-day Wednesday. Driver to call consignee 1 hour ahead.",
                "Lift gate required at delivery. Consignee will provide forklift at receiving dock.",
            ),
            max_tokens=110,
        )
        footer.append(f"<b>Special handling:</b> {instructions}")

    pdf = _pdf_table_doc(
        title=_rotate(
            f"Bill of Lading — {bol_no}", f"BOL {bol_no}", f"Carrier Bill of Lading — {bol_no}"
        ),
        header_pairs=pairs,
        table_header=["Description", "Pkgs", "Weight", "Freight class"],
        table_rows=rows,
        footer_lines=footer,
        pdf_title="Bill of lading",
    )
    return pdf, {
        "bol_no": bol_no,
        "shipper": shipper,
        "consignee": consignee,
        "line_count": len(rows),
    }


def _gen_pdf_customs(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a customs declaration.

    Variations: 2-8 declared items, AI-drafted purpose-of-shipment
    statement (50% of the time), wider incoterm pool, optional
    duty-paid summary.
    """
    declarant = fkr.company()
    decl_no = f"CD-{datetime.now().strftime('%Y%m%d')}-{random.randint(10000, 99999)}"
    incoterm = random.choice(_ctx(industry, "incoterms", ["FOB"]))
    origin = fkr.country()
    destination = fkr.country()
    rows = []
    total_value = 0.0
    line_count = random.randint(2, 8)
    for _ in range(line_count):
        item = fkr.bs().title()
        hts = f"{random.randint(1000, 9999)}.{random.randint(10, 99)}.{random.randint(1000, 9999)}"
        qty = random.randint(10, 1000)
        value = round(random.uniform(100, 25000), 2)
        total_value += value
        rows.append([item, hts, str(qty), f"${value:,.2f}"])
    rows.append(["", "", "Total declared value", f"${total_value:,.2f}"])

    pairs = [
        ("Declarant", declarant),
        ("Country of origin", origin),
        ("Country of destination", destination),
        ("Incoterm", incoterm),
        ("Declaration #", decl_no),
    ]
    if _maybe_section(0.5):
        purpose = _maybe_ai(
            ai_client,
            f"Write a 1-sentence 'Purpose of shipment' declaration for a customs filing from {origin} to {destination}, declared value ${total_value:.2f}. Tone: official, customs-formal.",
            fallback=_rotate(
                "Commercial sale to end customer",
                "Sample shipment for customer evaluation",
                "Replacement parts under warranty",
                "Inter-company transfer to subsidiary",
            ),
            max_tokens=60,
        )
        pairs.append(("Purpose", purpose))
    if _maybe_section(0.4):
        duty_pct = random.choice([0.0, 2.5, 4.5, 7.5, 12.0])
        pairs.append(
            ("Estimated duty", f"{duty_pct}% (${round(total_value * duty_pct / 100, 2):,.2f})")
        )

    pdf = _pdf_table_doc(
        title=_rotate(
            f"Customs Declaration — {decl_no}",
            f"Export Declaration — {decl_no}",
            f"Commercial Invoice for Customs — {decl_no}",
        ),
        header_pairs=pairs,
        table_header=["Description", "HTS code", "Qty", "Value (USD)"],
        table_rows=rows,
        footer_lines=[
            "<i>"
            + _rotate(
                "I declare that the information given is true and complete. False declarations are subject to penalties under applicable customs law.",
                "Declaration accurate to the best of declarant's knowledge. Subject to verification by competent authorities.",
                "All goods listed have been correctly classified per the Harmonized Tariff Schedule.",
            )
            + "</i>",
        ],
        pdf_title="Customs declaration",
    )
    return pdf, {
        "decl_no": decl_no,
        "incoterm": incoterm,
        "origin": origin,
        "destination": destination,
        "line_count": line_count,
        "total_usd": round(total_value, 2),
    }


# Insurance ──


def _gen_pdf_underwriting_report(
    industry: str, fkr: Any, ai_client: Any | None
) -> tuple[bytes, dict]:
    """Synthesize an underwriting report.

    Variations: 4-8 underwriting factors picked from the wider pool,
    AI-drafted risk-assessment narrative + decision rationale (60%
    of the time), variable claim history.
    """
    applicant = fkr.name()
    policy = random.choice(_ctx(industry, "policy_types", ["Auto"]))
    underwriter = fkr.name()
    risk = random.choice(
        ["Preferred Plus", "Preferred", "Standard", "Standard", "Substandard", "High-risk"]
    )
    premium = round(random.uniform(300, 12500), 2)
    decision = random.choice(
        [
            "Approved",
            "Approved",
            "Approved with conditions",
            "Approved with conditions",
            "Declined",
            "Pending additional information",
        ]
    )
    claim_count = random.randint(0, 8)

    full_factors = [
        ("Loss history (5-year)", f"{claim_count} claims"),
        ("Coverage limit", f"${random.choice([100000, 250000, 500000, 1000000, 2000000]):,}"),
        ("Deductible", f"${random.choice([250, 500, 1000, 2500, 5000]):,}"),
        ("Risk classification", risk),
        ("Term", random.choice(["6 months", "12 months", "12 months", "24 months"])),
        ("Credit-based score", random.choice(["Excellent", "Good", "Fair", "Poor"])),
        (
            "Prior insurance",
            random.choice(["Continuous 5+ yrs", "Continuous 1-5 yrs", "Lapse in past 12 months"]),
        ),
        ("Geographic risk zone", random.choice(["Low", "Moderate", "Elevated", "High"])),
        ("Occupation class", random.choice(["Professional", "Standard", "Manual labour", "Other"])),
    ]
    factors = random.sample(full_factors, k=random.randint(4, 8))

    footer = [
        "<i>"
        + _rotate(
            "This report is for internal use. Final policy terms are subject to issuance and applicable state filings.",
            "Underwriting decision contingent on satisfactory completion of all conditions.",
            "Filed for internal underwriting record; not for distribution outside the underwriting department.",
        )
        + "</i>"
    ]
    if _maybe_section(0.6):
        rationale = _maybe_ai(
            ai_client,
            f"Write a 2-sentence underwriting decision rationale for {applicant}, {policy}, decision={decision}, risk={risk}, {claim_count} prior claims. Tone: insurance-formal, factual.",
            fallback=_rotate(
                f"Risk profile is consistent with the {risk} class given the loss-history pattern and other factors. Decision: {decision}.",
                f"Application reviewed against current underwriting guidelines; {claim_count} prior claims align with {risk} pricing tier. {decision}.",
                f"Underwriting analysis supports the {decision.lower()} disposition based on the totality of risk factors presented.",
            ),
            max_tokens=140,
        )
        footer.append(f"<b>Decision rationale:</b> {rationale}")

    pdf = _pdf_table_doc(
        title=_rotate(
            "Underwriting Report",
            "Underwriting Decision Memo",
            "Risk Assessment & Underwriting Report",
        ),
        header_pairs=[
            ("Applicant", applicant),
            ("Policy type", policy),
            ("Underwriter", underwriter),
            ("Decision", decision),
            ("Quoted premium", f"${premium:,.2f}"),
        ],
        table_header=["Underwriting factor", "Value"],
        table_rows=[list(r) for r in factors],
        footer_lines=footer,
        pdf_title="Underwriting report",
    )
    return pdf, {
        "applicant": applicant,
        "policy": policy,
        "decision": decision,
        "premium_usd": premium,
        "risk": risk,
        "factor_count": len(factors),
        "claim_count": claim_count,
    }


def _gen_docx_endorsement(industry: str, fkr: Any, ai_client: Any | None) -> tuple[bytes, dict]:
    """Synthesize a policy endorsement letter.

    Variations: AI-drafted endorsement-description paragraph,
    optional premium adjustment paragraph, rotated closing.
    """
    policy_no = f"POL-{random.randint(1000000, 9999999)}"
    policyholder = fkr.name()
    code = random.choice(_ctx(industry, "endorsement_codes", ["END-001"]))
    effective = datetime.now().date().isoformat()
    has_premium_adj = _maybe_section(0.5)
    premium_delta = round(random.uniform(-200, 800), 2) if has_premium_adj else 0.0

    paragraphs = [
        f"Dear {policyholder.split()[0]}, this letter confirms the addition of endorsement {code} to your policy effective {effective}.",
        _maybe_ai(
            ai_client,
            f"Write a 2-sentence description of insurance endorsement {code} added to policy {policy_no}. Tone: insurance-formal, customer-facing.",
            fallback="The endorsement modifies your existing coverage as described in the attached endorsement form. Please review the changes and retain this document with your policy materials.",
            max_tokens=140,
        ),
    ]
    if has_premium_adj:
        if premium_delta > 0:
            paragraphs.append(
                f"This endorsement increases your annual premium by ${premium_delta:.2f}; the adjustment will appear on your next billing statement."
            )
        else:
            paragraphs.append(
                f"This endorsement reduces your annual premium by ${abs(premium_delta):.2f}; a credit will appear on your next billing statement."
            )
    else:
        paragraphs.append(
            "There is no premium adjustment associated with this endorsement at this time. If your coverage needs change, please contact your agent."
        )

    paragraphs.append(
        _rotate(
            "Thank you for your continued business.",
            "We appreciate the trust you place in us as your insurance partner.",
            "If you have any questions, your agent is available to walk through the changes.",
        )
    )

    docx = _docx_letter_doc(
        sender_name=fkr.company() + " Insurance",
        sender_address=fkr.address(),
        recipient_name=policyholder,
        recipient_address=fkr.address(),
        subject=f"Policy endorsement — {policy_no}",
        body_paragraphs=paragraphs,
        closing=_rotate("Sincerely,", "Best regards,", "Kind regards,", "Yours faithfully,"),
    )
    return docx, {
        "policy_no": policy_no,
        "endorsement_code": code,
        "effective": effective,
        "premium_delta": premium_delta,
    }


# ── Top-level orchestrator ────────────────────────────────────────


def _ensure_volume(client: WorkspaceClient, warehouse_id: str, vol_fqn: str) -> None:
    """Idempotent CREATE VOLUME IF NOT EXISTS. Same pattern the
    convert smoke endpoint uses; mirrored here so the Documents
    generator "just works" against a fresh schema with no manual
    setup."""
    execute_sql(client, warehouse_id, f"CREATE VOLUME IF NOT EXISTS {vol_fqn}")


def _ensure_catalog_table(
    client: WorkspaceClient,
    warehouse_id: str,
    fqn: str,
    *,
    direct: bool,
) -> None:
    """Create-or-replace the catalog/direct table. ``direct=True``
    emits a ``content BINARY`` column for inline file bytes; ``False``
    emits a ``file_path STRING`` column that points at the Volume URI.

    Drop-and-create behaviour matches the structured generator's
    ``CREATE OR REPLACE`` semantics — generating a demo catalog should
    never silently merge into an existing table from a previous run.
    """
    if direct:
        sql = f"""
        CREATE OR REPLACE TABLE {fqn} (
            file_id          STRING,
            doc_type         STRING,
            file_extension   STRING,
            size_bytes       BIGINT,
            industry         STRING,
            generated_at     TIMESTAMP,
            content_summary  STRING,
            page_count       BIGINT,
            content          BINARY,
            metadata_json    STRING
        ) USING delta
        """
    else:
        sql = f"""
        CREATE OR REPLACE TABLE {fqn} (
            file_path        STRING,
            doc_type         STRING,
            file_extension   STRING,
            size_bytes       BIGINT,
            industry         STRING,
            generated_at     TIMESTAMP,
            content_summary  STRING,
            content_full     STRING,
            page_count       BIGINT,
            metadata_json    STRING
        ) USING delta
        """
    execute_sql(client, warehouse_id, sql)


def _build_content_full(meta: dict) -> str:
    """Flat textual projection of a document for the indexed catalog.

    Documents are binary (PDF / DOCX / PPTX / XLSX / EML) so the
    inline bytes can't be searched from SQL. We synthesise a queryable
    text blob from the string-shaped fields each generator already
    stashes on its return-meta — claim diagnoses, invoice line items,
    letter body paragraphs, report sections, deck slides, etc.

    Heuristic, not a full text extraction: catches string and
    list-of-string fields, skips numeric / structured fields. Good
    enough for `WHERE content_full LIKE '%billing%'` style RAG
    queries without forcing every generator to return identical
    metadata shapes.
    """
    parts: list[str] = []

    def _add(value: Any) -> None:
        if value is None:
            return
        if isinstance(value, str):
            text = value.strip()
            # Skip short identifier-shaped fields (case_id, "PA chest")
            # that don't add search value.
            if len(text) >= 8 and not text.replace("-", "").replace("_", "").isalnum():
                parts.append(text)
            elif len(text) >= 12:  # longer alnum strings still useful
                parts.append(text)
        elif isinstance(value, (list, tuple)):
            for item in value:
                _add(item)

    for value in meta.values():
        _add(value)
    return "\n\n".join(parts)


def _sql_str(s: str | None) -> str:
    """Single-quote escape a string for inline INSERT VALUES. Doubles
    embedded apostrophes; returns ``NULL`` for None."""
    if s is None:
        return "NULL"
    return "'" + s.replace("'", "''") + "'"


def generate_documents(
    client: WorkspaceClient,
    warehouse_id: str,
    config: dict,
    progress: dict | None = None,
    stop_check: Callable[[], bool] | None = None,
) -> dict:
    """Top-level orchestrator. Generates and lands the requested doc
    mix, returning a summary dict the API router serializes back to
    the UI.

    ``config`` keys:
        catalog, schema, volume:    str (catalog/schema always required;
                                          volume required when destination
                                          is `volume` or `volume_with_catalog`)
        types:                      list[str] — registry keys
        counts:                     dict[str, int] — per-type count
        industry:                   str (default "healthcare")
        destination:                Literal["volume", "volume_with_catalog",
                                            "direct_table"] (default
                                            "volume_with_catalog")
        realistic_content:          bool (default False — AI mode)
        faker_locale:               str (default "en_US")
        faker_seed:                 int | None (default None)

    ``progress`` (mutated in place per file):
        {files_written, total_bytes, current_type, current_path,
         per_type: dict[str, int]}

    ``stop_check`` (returns True to stop early — the JobManager wires
    this to the operator's "Stop" button).
    """
    if not DOCUMENTS_AVAILABLE:
        raise RuntimeError(_UNAVAILABLE_REASON or "documents extra not installed")

    progress = progress if progress is not None else {}
    stopped = stop_check or (lambda: False)

    catalog = config["catalog"]
    schema = config["schema"]
    types = config.get("types") or []
    counts = config.get("counts") or {}
    industry = config.get("industry", "healthcare")
    destination = config.get("destination", "volume_with_catalog")

    if destination not in ("volume", "volume_with_catalog", "direct_table"):
        raise ValueError(f"Unknown destination: {destination!r}")
    if not types:
        raise ValueError("'types' must contain at least one document type")
    unknown = [t for t in types if t not in DOCUMENT_TYPES]
    if unknown:
        raise ValueError(f"Unknown document types: {unknown}. Known: {sorted(DOCUMENT_TYPES)}")

    # Faker pool — reused across every generator call to keep names
    # internally consistent (Alice's claim form may be referenced in
    # Alice's email; same Faker instance, same seed, same Alice).
    from faker import Faker

    fkr = Faker(locale=config.get("faker_locale", "en_US"))
    if config.get("faker_seed") is not None:
        fkr.seed_instance(int(config["faker_seed"]))

    # AI client — opt-in via realistic_content. ``build_drafter``
    # returns None when AI mode is disabled or no backend is
    # configured; generators handle the None case gracefully.
    from src.ai_drafter import adapter_summary, build_drafter

    ai_client = build_drafter(config, sdk_client=client)

    # Volume + table setup (skipped per-destination)
    volume_path: str | None = None
    table_fqn: str | None = None
    if destination in ("volume", "volume_with_catalog"):
        volume = config.get("volume") or "demo_unstructured"
        vol_fqn = f"{catalog}.{schema}.{volume}"
        _ensure_volume(client, warehouse_id, vol_fqn)
        volume_path = f"/Volumes/{catalog}/{schema}/{volume}/documents"

    # Custom table name (optional) — overrides the default
    # demo_documents_catalog / demo_documents when an operator wants
    # distinct table namespaces across runs.
    custom_table = (config.get("table_name") or "").strip()
    if destination == "volume_with_catalog":
        table_fqn = f"{catalog}.{schema}.{custom_table or 'demo_documents_catalog'}"
        _ensure_catalog_table(client, warehouse_id, table_fqn, direct=False)
    elif destination == "direct_table":
        table_fqn = f"{catalog}.{schema}.{custom_table or 'demo_documents'}"
        _ensure_catalog_table(client, warehouse_id, table_fqn, direct=True)

    # Initialise progress
    progress.setdefault("files_written", 0)
    progress.setdefault("total_bytes", 0)
    progress.setdefault("per_type", {t: 0 for t in types})
    progress.setdefault("destination", destination)

    # Per-row INSERT batches — accumulate up to 50 rows then flush so
    # we don't issue one INSERT per file (slow) but also don't risk a
    # mid-run crash dropping more than ~50 row's worth of metadata.
    pending_rows: list[str] = []
    BATCH_SIZE = 50

    def _flush_pending() -> None:
        nonlocal pending_rows
        if not pending_rows or table_fqn is None:
            return
        if destination == "volume_with_catalog":
            cols = (
                "file_path",
                "doc_type",
                "file_extension",
                "size_bytes",
                "industry",
                "generated_at",
                "content_summary",
                "content_full",
                "page_count",
                "metadata_json",
            )
        else:  # direct_table
            cols = (
                "file_id",
                "doc_type",
                "file_extension",
                "size_bytes",
                "industry",
                "generated_at",
                "content_summary",
                "page_count",
                "content",
                "metadata_json",
            )
        sql = f"INSERT INTO {table_fqn} ({', '.join(cols)}) VALUES {', '.join(pending_rows)}"
        execute_sql(client, warehouse_id, sql)
        pending_rows = []

    started_at = datetime.now(timezone.utc)

    for type_id in types:
        if stopped():
            break
        n = int(counts.get(type_id, 10))
        type_def = DOCUMENT_TYPES[type_id]
        gen_fn_name = type_def["gen_fn"]
        gen_fn = globals().get(gen_fn_name)
        if gen_fn is None:
            logger.error(f"Generator not found: {gen_fn_name}")
            continue
        progress["current_type"] = type_id

        for seq in range(n):
            if stopped():
                break
            try:
                file_bytes, meta = gen_fn(industry, fkr, ai_client)
            except Exception as e:
                logger.error(f"  ✗ {type_id} #{seq}: {e}")
                continue
            file_id = uuid.uuid4().hex
            ext = type_def["extension"]
            file_name = f"{type_id}_{file_id}.{ext}"

            # ── Write to Volume ──
            current_path: str | None = None
            if volume_path is not None:
                current_path = f"{volume_path}/{type_id}/{file_name}"
                client.files.upload(
                    file_path=current_path,
                    contents=io.BytesIO(file_bytes),
                    overwrite=True,
                )

            # ── Build the row to INSERT ──
            content_summary = meta.get("content_summary") or _build_summary(type_id, meta)
            content_full = _build_content_full(meta)
            page_count = int(meta.get("page_count") or 0)
            metadata_json = json.dumps(meta, default=str)

            if destination == "volume_with_catalog" and current_path:
                row = (
                    f"({_sql_str(current_path)}, "
                    f"{_sql_str(type_id)}, "
                    f"{_sql_str(ext)}, "
                    f"{len(file_bytes)}, "
                    f"{_sql_str(industry)}, "
                    f"current_timestamp(), "
                    f"{_sql_str(content_summary)}, "
                    f"{_sql_str(content_full)}, "
                    f"{page_count}, "
                    f"{_sql_str(metadata_json)})"
                )
                pending_rows.append(row)
            elif destination == "direct_table":
                # Inline bytes via unhex(hex(...)) is impractical for
                # multi-MB files — use the SDK's parameterized SQL
                # path instead. For v1 simplicity we INSERT each
                # direct-table row individually using a parameterised
                # statement (slower but correct). Batched binary
                # INSERTs are a v2 optimisation.
                _insert_direct_row(
                    client,
                    warehouse_id,
                    table_fqn or "",
                    file_id=file_id,
                    doc_type=type_id,
                    file_extension=ext,
                    size_bytes=len(file_bytes),
                    industry=industry,
                    content_summary=content_summary,
                    page_count=page_count,
                    content=file_bytes,
                    metadata_json=metadata_json,
                )

            # ── Update progress ──
            progress["files_written"] = progress.get("files_written", 0) + 1
            progress["total_bytes"] = progress.get("total_bytes", 0) + len(file_bytes)
            progress["per_type"][type_id] = progress["per_type"].get(type_id, 0) + 1
            if current_path:
                progress["current_path"] = current_path

            # Flush per-type or every BATCH_SIZE rows
            if destination == "volume_with_catalog" and len(pending_rows) >= BATCH_SIZE:
                _flush_pending()

    _flush_pending()

    completed_at = datetime.now(timezone.utc)
    duration_ms = int((completed_at - started_at).total_seconds() * 1000)
    result: dict[str, Any] = {
        "status": "completed",
        "files_written": progress["files_written"],
        "total_bytes": progress["total_bytes"],
        "per_type": progress["per_type"],
        "destination": destination,
        "volume_path": volume_path,
        "table_fqn": table_fqn,
        "duration_ms": duration_ms,
        "started_at": started_at.isoformat(timespec="seconds"),
        "completed_at": completed_at.isoformat(timespec="seconds"),
    }
    result.update(adapter_summary(ai_client))
    return result


def _insert_direct_row(
    client: WorkspaceClient,
    warehouse_id: str,
    table_fqn: str,
    *,
    file_id: str,
    doc_type: str,
    file_extension: str,
    size_bytes: int,
    industry: str,
    content_summary: str,
    page_count: int,
    content: bytes,
    metadata_json: str,
) -> None:
    """Insert one row into the direct table with the binary `content`.

    Direct-table inserts are slower than Volume + catalog because of
    the per-row binary payload. v1 ships them per-row; batched bulk
    inserts (Spark Connect / temp staging) are a v2 optimisation.
    """
    # Use unhex() of the hex-encoded payload — the most portable way
    # to embed BINARY in a SQL string across DBSQL versions.
    hex_content = content.hex()
    sql = (
        f"INSERT INTO {table_fqn} "
        f"(file_id, doc_type, file_extension, size_bytes, industry, "
        f"generated_at, content_summary, page_count, content, metadata_json) "
        f"VALUES ("
        f"{_sql_str(file_id)}, "
        f"{_sql_str(doc_type)}, "
        f"{_sql_str(file_extension)}, "
        f"{size_bytes}, "
        f"{_sql_str(industry)}, "
        f"current_timestamp(), "
        f"{_sql_str(content_summary)}, "
        f"{page_count}, "
        f"unhex('{hex_content}'), "
        f"{_sql_str(metadata_json)})"
    )
    execute_sql(client, warehouse_id, sql)


def _build_summary(type_id: str, meta: dict) -> str:
    """Build a human-readable one-line summary from the per-type
    metadata dict. Used as `content_summary` in the catalog table so
    SQL queries like ``WHERE content_summary LIKE '%hypertension%'``
    work without parsing the metadata_json blob."""
    if type_id == "pdf_claim":
        return f"Claim {meta.get('claim_id')} for {meta.get('patient_name')} (${meta.get('total_charges', 0):.2f})"
    if type_id == "pdf_invoice":
        return (
            f"Invoice {meta.get('invoice_id')} from {meta.get('vendor')} to {meta.get('customer')}"
        )
    if type_id == "pdf_contract":
        return f"Services Agreement between {meta.get('party_a')} and {meta.get('party_b')}"
    if type_id == "docx_letter":
        return f"Letter from {meta.get('sender')} to {meta.get('recipient')} re: {meta.get('department')}"
    if type_id == "docx_report":
        return f"{meta.get('quarter')} {meta.get('department')} report by {meta.get('author')}"
    if type_id == "pptx_deck":
        return f"{meta.get('company')} pitch to {meta.get('audience')} ({meta.get('slide_count')} slides)"
    if type_id == "xlsx_budget":
        return f"FY{meta.get('fiscal_year')} {meta.get('department')} budget (${meta.get('annual_total', 0):,})"
    if type_id == "xlsx_inventory":
        return f"Inventory list — {meta.get('item_count')} items across {meta.get('category_count')} categories"
    if type_id == "eml_message":
        return f"Email from {meta.get('sender')} to {meta.get('recipient')}: {meta.get('subject')}"
    return type_id


# ── Preview (pure arithmetic) ─────────────────────────────────────


def preview_documents(config: dict) -> dict:
    """Return per-type / total estimates without going near the
    warehouse. Same shape as the structured generator's preview —
    operators can see "approximately N MB total, ~M minutes" before
    they click Generate.

    ``config`` shape mirrors ``generate_documents`` but only `types`
    and `counts` are read here.
    """
    types = config.get("types") or []
    counts = config.get("counts") or {}
    per_type = []
    total_files = 0
    total_bytes = 0
    total_seconds = 0.0
    unknown: list[str] = []
    for t in types:
        if t not in DOCUMENT_TYPES:
            unknown.append(t)
            continue
        n = int(counts.get(t, 10))
        bytes_each = _AVG_BYTES_PER_TYPE.get(t, 5_000)
        per_sec = _GEN_PER_SECOND_PER_TYPE.get(t, 100)
        per_type.append(
            {
                "type": t,
                "category": DOCUMENT_TYPES[t]["category"],
                "label": DOCUMENT_TYPES[t]["label"],
                "count": n,
                "estimated_bytes": n * bytes_each,
                "estimated_seconds": round(n / per_sec, 1),
            }
        )
        total_files += n
        total_bytes += n * bytes_each
        total_seconds += n / per_sec
    return {
        "per_type": per_type,
        "total_files": total_files,
        "total_bytes": total_bytes,
        "estimated_seconds": round(total_seconds, 1),
        "unknown_types": unknown,
    }
